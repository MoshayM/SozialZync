from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── colour palette ─────────────────────────────────────────────────────────────
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK    = RGBColor(0x1A, 0x1A, 0x2E)
C_PURPLE  = RGBColor(0x7C, 0x3A, 0xED)
C_PURPLE2 = RGBColor(0x6D, 0x28, 0xD9)
C_AMBER   = RGBColor(0xD9, 0x77, 0x06)
C_AMBER2  = RGBColor(0xB4, 0x5A, 0x09)
C_NAVY    = RGBColor(0x1A, 0x1A, 0x2E)
C_NAVY2   = RGBColor(0x16, 0x21, 0x3E)
C_TEAL    = RGBColor(0x08, 0x91, 0xB2)
C_TEAL2   = RGBColor(0x06, 0x6A, 0x85)
C_GRAY    = RGBColor(0xF1, 0xF0, 0xFA)
C_LGRAY   = RGBColor(0xE5, 0xE5, 0xF0)

blank_layout = prs.slide_layouts[6]  # completely blank

def rgb(r, g, b):
    return RGBColor(r, g, b)

def set_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    return shape

def add_text_box(slide, text, left, top, width, height,
                 font_size=16, bold=False, color=C_WHITE,
                 align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.name = 'Calibri'
    return txBox

def add_slide_number(slide, num, total, color=C_WHITE):
    add_text_box(slide, f"{num} / {total}",
                 left=11.8, top=7.1, width=1.4, height=0.3,
                 font_size=10, color=color, align=PP_ALIGN.RIGHT)

def make_title_slide(slide, title, subtitle, tagline, url, date_str, slide_num, total):
    # Background gradient effect using two rectangles
    set_bg(slide, rgb(0x1A, 0x0D, 0x3E))
    add_rect(slide, 0, 0, 13.33, 7.5, rgb(0x1A, 0x0D, 0x3E))
    add_rect(slide, 0, 0, 13.33, 4.0, rgb(0x2D, 0x1B, 0x69))

    # Accent bar
    add_rect(slide, 1.2, 2.9, 0.08, 2.0, C_PURPLE)

    # Title
    add_text_box(slide, title,
                 left=1.5, top=1.3, width=10.5, height=1.6,
                 font_size=52, bold=True, color=C_WHITE)

    # Subtitle
    add_text_box(slide, subtitle,
                 left=1.5, top=3.05, width=10.0, height=0.6,
                 font_size=22, bold=False, color=rgb(0xC4, 0xB5, 0xFD))

    # Tagline
    add_text_box(slide, tagline,
                 left=1.5, top=3.75, width=10.0, height=0.5,
                 font_size=14, bold=False, color=rgb(0xA7, 0x8B, 0xFA))

    # URL
    add_text_box(slide, f"  {url}",
                 left=1.5, top=5.6, width=8.0, height=0.4,
                 font_size=13, bold=False, color=rgb(0xC4, 0xB5, 0xFD))

    # Date
    add_text_box(slide, date_str,
                 left=1.5, top=6.15, width=8.0, height=0.35,
                 font_size=11, bold=False, color=rgb(0x9C, 0x8B, 0xB9))

    add_slide_number(slide, slide_num, total, color=rgb(0x9C, 0x8B, 0xB9))

def make_section_divider(slide, section_label, section_title, accent_color, slide_num, total):
    set_bg(slide, accent_color)
    # Decorative rectangle
    add_rect(slide, 0, 0, 0.5, 7.5, rgb(
        min(accent_color.red + 30, 255),
        min(accent_color.green + 30, 255),
        min(accent_color.blue + 30, 255)
    ))
    # Section label (small)
    add_text_box(slide, section_label,
                 left=1.2, top=2.4, width=11.0, height=0.5,
                 font_size=14, bold=False, color=rgb(0xFF, 0xFF, 0xFF),
                 align=PP_ALIGN.LEFT)
    # Section title (large)
    add_text_box(slide, section_title,
                 left=1.2, top=2.9, width=11.0, height=1.5,
                 font_size=40, bold=True, color=C_WHITE,
                 align=PP_ALIGN.LEFT)
    add_slide_number(slide, slide_num, total)

def make_content_slide(slide, title, bullets, accent_color, slide_num, total, header_height=0.75):
    set_bg(slide, C_GRAY)
    # Header bar
    add_rect(slide, 0, 0, 13.33, header_height, accent_color)
    # Title in header
    add_text_box(slide, title,
                 left=0.35, top=0.05, width=12.5, height=header_height - 0.05,
                 font_size=22, bold=True, color=C_WHITE)

    # Content area — white card
    add_rect(slide, 0.3, header_height + 0.15, 12.73, 7.5 - header_height - 0.35,
             C_WHITE)

    # Bullet text
    content_top = header_height + 0.3
    content_h   = 7.5 - header_height - 0.55

    txBox = slide.shapes.add_textbox(
        Inches(0.55), Inches(content_top), Inches(12.2), Inches(content_h)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    for bullet in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        indent = bullet.startswith('   ')
        text   = bullet.lstrip()

        p.space_before = Pt(4 if not indent else 2)
        p.level = 1 if indent else 0

        run = p.add_run()
        run.text = text
        run.font.name = 'Calibri'
        run.font.size = Pt(13 if not indent else 11.5)
        run.font.bold = not indent and text.startswith('●')
        run.font.color.rgb = C_DARK if not indent else rgb(0x44, 0x44, 0x66)

    add_slide_number(slide, slide_num, total, color=rgb(0x88, 0x88, 0xAA))

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE DEFINITIONS
# ══════════════════════════════════════════════════════════════════════════════

TOTAL = 29
slide_num = 0

# ── Slide 1 — Title ────────────────────────────────────────────────────────────
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
make_title_slide(sl,
    title="AI CreatorForce",
    subtitle="AI-Powered YouTube Content Operating System",
    tagline="From idea to published video — with a full AI content team behind every creator.",
    url="sozialzync.vercel.app",
    date_str="July 2026",
    slide_num=slide_num, total=TOTAL)

# ── Slide 2 — What is AI CreatorForce ─────────────────────────────────────────
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
make_content_slide(sl,
    title="What is AI CreatorForce?",
    bullets=[
        "● Production-grade SaaS platform — a complete AI content operating system for YouTube creators.",
        "● Automates the entire content lifecycle: Research → Script → Fact-check → Compliance → SEO → Publish.",
        "● Every output must pass a hard compliance gate before it can be published — no bypasses.",
        "● Human-in-the-loop on publish: the AI drafts, the creator approves before anything reaches YouTube.",
        "● Supports Free, Pro, Enterprise, and Agency plan tiers; SUPER_ADMIN/OWNER have full access to everything.",
        "● Built on Next.js (frontend) + NestJS (backend) + 10+ specialized AI agents + BullMQ async pipeline.",
        "● Not a spam generator — every feature enforces original, monetizable, copyright-safe content.",
    ],
    accent_color=C_PURPLE, slide_num=slide_num, total=TOTAL)

# ══════════════════════════════════════════════════════════════════════════════
#  SECTION 2 — CREATOR PERSPECTIVE
# ══════════════════════════════════════════════════════════════════════════════

# ── Slide 3 — Your AI Content Team ────────────────────────────────────────────
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
make_content_slide(sl,
    title="FOR THE CREATOR  |  Your AI Content Team",
    bullets=[
        "● ResearchAgent — discovers trending topics, pulls sources, builds a fact-sourced research brief.",
        "● ScriptAgent — writes hook, body, and CTA sections tailored to your channel voice and niche.",
        "● FactCheckAgent — verifies every factual claim against ResearchAgent sources; rejects hallucinations.",
        "● ComplianceAgent — scores content 0–100; blocks anything below 70 or with BLOCK-severity flags.",
        "● SEOAgent — optimizes title variants, tags, and keyword targeting for YouTube search.",
        "● MetadataAgent — generates video title, description, chapters, and thumbnail brief automatically.",
        "● ThumbnailAgent — generates AI thumbnail candidates; you pick the winner.",
        "● CopilotAgent — always-on chat assistant for questions, edits, and content strategy during creation.",
    ],
    accent_color=C_PURPLE, slide_num=slide_num, total=TOTAL)

# ── Slide 4 — The Content Pipeline ────────────────────────────────────────────
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
make_content_slide(sl,
    title="FOR THE CREATOR  |  The Content Pipeline (Long-Form)",
    bullets=[
        "● Step 1 — Research:    ResearchAgent sources topic data, competitor intel, trending angles.",
        "● Step 2 — Script:      ScriptAgent writes a structured script with hook, body, CTA.",
        "● Step 3 — Fact Check:  FactCheckAgent traces every claim to a ResearchAgent source.",
        "● Step 4 — Compliance:  ComplianceAgent hard gate — score ≥ 70, no BLOCK flags → pass.",
        "● Step 5 — Metadata:    MetadataAgent generates title, description, chapters, thumbnail brief.",
        "● Step 6 — SEO:         SEOAgent optimizes title variants, tags, and keyword strategy.",
        "● Step 7 — You Review:  Approvals UI shows the full package for your sign-off.",
        "● Step 8 — Publish:     PublishingService uploads to YouTube only after your approval.",
        "   ○ Every step runs as an async BullMQ job — nothing blocks the UI while AI works.",
        "   ○ Failed steps route to QualityControlAgent for retry; you get notified of blockers.",
    ],
    accent_color=C_PURPLE, slide_num=slide_num, total=TOTAL)

# ── Slide 5 — Plan Tiers ───────────────────────────────────────────────────────
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
make_content_slide(sl,
    title="FOR THE CREATOR  |  Plan Tiers",
    bullets=[
        "● FREE — 3 projects, 5 outputs per project, 10 Copilot queries/day, 10 Shorts edits/month.",
        "   ○ Access to basic AI tools; Compliance gate always active.",
        "   ○ Great for testing the platform and getting your first AI-assisted video done.",
        "● PRO — Credits-based access. Top up wallet credits and unlock all 15 AI agents.",
        "   ○ Unlimited projects, unlimited publishing, full SEO suite, Shorts Studio.",
        "   ○ No fixed monthly subscription: pay per generation with credit lots.",
        "   ○ Credits never expire faster than the lot expiry date; multiple lot types supported.",
        "● ENTERPRISE — Team workspaces, shared org wallet, multi-channel management.",
        "   ○ Org admin controls budget caps, team allocations, and member roles.",
        "   ○ White-label dashboard, SLA, priority support, dedicated account manager.",
        "● SUPER_ADMIN / OWNER — Full platform access; bypasses all plan gates automatically.",
    ],
    accent_color=C_PURPLE, slide_num=slide_num, total=TOTAL)

# ── Slide 6 — Credit & Wallet System ──────────────────────────────────────────
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
make_content_slide(sl,
    title="FOR THE CREATOR  |  Credit & Wallet System",
    bullets=[
        "● Pay-as-you-go — no forced subscription. Buy credits when you need them.",
        "● Credit lot types: Trial (free on sign-up), Purchased, Bonus, Referral rewards.",
        "   ○ Lots are consumed in expiry order (oldest first) — FIFO to protect your credits.",
        "● Wallet balance is live: queried via API, shown in the header and wallet page.",
        "● Credits below 500 trigger a 'running low' warning badge in the UI.",
        "● Referral program: share your referral code → earn bonus credits when friends join.",
        "● Trial grant on sign-up: try AI agents immediately without entering payment info.",
        "● Credit reservation system: credits are held during a job and settled on completion.",
        "   ○ If a job fails, the reservation is released — you are never charged for failed runs.",
        "● Stripe integration for purchasing credit packs or upgrading to Pro/Enterprise subscription.",
    ],
    accent_color=C_PURPLE, slide_num=slide_num, total=TOTAL)

# ── Slide 7 — Publishing Workflow ─────────────────────────────────────────────
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
make_content_slide(sl,
    title="FOR THE CREATOR  |  Publishing Workflow",
    bullets=[
        "● Human approval is mandatory — the platform never publishes to YouTube without your sign-off.",
        "● Approval UI shows: AI script, compliance score, SEO metadata, thumbnail, chapter list.",
        "● You can edit any field before approving — the final call is always yours.",
        "● Once approved, PublishingService calls YouTube Data API with the full metadata package.",
        "● Scheduling: set a future publish date and time; the job fires automatically at the scheduled slot.",
        "   ○ Auto-publish (scheduled) is only available if the item has already passed compliance.",
        "● YouTube OAuth is required: connect your channel in Settings → Channels before publishing.",
        "● Chapter timestamps are written into the video description automatically.",
        "● Post-publish: the platform polls YouTube Analytics to track views, watch time, and CTR.",
    ],
    accent_color=C_PURPLE, slide_num=slide_num, total=TOTAL)

# ── Slide 8 — Shorts Studio ───────────────────────────────────────────────────
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
make_content_slide(sl,
    title="FOR THE CREATOR  |  Shorts Studio",
    bullets=[
        "● Turn your existing long-form videos into YouTube Shorts — no re-recording needed.",
        "● Import flow: select a channel → pick videos from your library → AI analyzes each one.",
        "   ○ Analysis: transcript ingestion, scene detection, topic segmentation, chapter detection.",
        "● Clip recommendations: AI ranks segments by virality potential, topic coherence, retention signals.",
        "● Timeline editor: drag-drop clip order, set trim points and transitions on a visual timeline.",
        "● AI editing assistant: type natural-language instructions ('make it punchier') → AI edits the timeline.",
        "● Thumbnail generator: AI creates Short thumbnail candidates; you approve the winner.",
        "● Social content factory: auto-generate Quote Cards, Carousels, Blog Posts, or Newsletters from the same content.",
        "● Same compliance gate applies: every Short passes ComplianceAgent before export or publish.",
        "● Chapter sync: confirmed chapters can be written back to your YouTube video description.",
    ],
    accent_color=C_PURPLE, slide_num=slide_num, total=TOTAL)

# ══════════════════════════════════════════════════════════════════════════════
#  SECTION 3 — ENTERPRISE OWNER PERSPECTIVE
# ══════════════════════════════════════════════════════════════════════════════

# ── Slide 9 — Enterprise Plan ─────────────────────────────────────────────────
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
make_content_slide(sl,
    title="FOR THE ENTERPRISE OWNER  |  Enterprise Plan — What You Get",
    bullets=[
        "● Unlimited projects, all 15 AI agents, full media pipeline, multi-channel management.",
        "● Team workspaces: create an Organization, invite team members, assign roles and budgets.",
        "● Shared org wallet: one billing account funds all team members' usage.",
        "● Budget controls: set hard monthly credit caps per team or per project — overage is blocked.",
        "● White-label dashboard: brand the platform UI with your company name and logo.",
        "● SLA guarantee: uptime commitment, priority issue resolution, dedicated support queue.",
        "● Dedicated account manager for onboarding, strategy, and feature requests.",
        "● Audit logs: every action (publish, approval, role change, billing event) is logged with timestamp.",
        "● Advanced analytics: BI module with cross-channel performance reports and growth tracking.",
        "● SUPER_ADMIN bypass: Owner/Developer account has unrestricted access to all features.",
    ],
    accent_color=C_AMBER, slide_num=slide_num, total=TOTAL)

# ── Slide 10 — Organization & Team Management ─────────────────────────────────
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
make_content_slide(sl,
    title="FOR THE ENTERPRISE OWNER  |  Organization & Team Management",
    bullets=[
        "● Organization model: one org can contain multiple teams, channels, and wallets.",
        "● Role hierarchy (Org level):",
        "   ○ ORG_ADMIN — full org control: billing, members, channels, budget allocation.",
        "   ○ BILLING_ADMIN — manages wallet top-ups, budget periods, credit allocation; no content access.",
        "   ○ TEAM_MANAGER — creates and manages teams, assigns EDITOR / REVIEWER / VIEWER roles.",
        "   ○ MEMBER — standard creator access within assigned team's scope and budget.",
        "● Team roles: OWNER → ADMIN → EDITOR → REVIEWER → VIEWER (progressively narrower access).",
        "● Budget periods: allocate a monthly credit budget per team; spend tracked in real time.",
        "   ○ Hard cap enforcement: credit reservation is rejected if it would exceed the team budget.",
        "● Shared wallet: org-level CreditLedger — all member spend deducted from the same pot.",
        "● Invite flow: email invite → OTP verification → role assignment → workspace access.",
    ],
    accent_color=C_AMBER, slide_num=slide_num, total=TOTAL)

# ── Slide 11 — Billing & Cost Control ─────────────────────────────────────────
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
make_content_slide(sl,
    title="FOR THE ENTERPRISE OWNER  |  Billing & Cost Control",
    bullets=[
        "● Shared org wallet: one Stripe payment method funds all team credit spending.",
        "● CreditLedger: append-only ledger — every credit movement is logged (grant/spend/reserve/settle/refund).",
        "● CreditLot buckets: Purchased, Trial, Bonus, Referral — each with its own expiry date.",
        "   ○ Lots consumed in FIFO expiry order: oldest expiring credits spent first.",
        "● BudgetPeriod: monthly credit cap per org or per team — set and forget.",
        "   ○ Approaching 80% budget threshold triggers a warning notification to ORG_ADMIN.",
        "● Per-project billing: each project tracks its credit spend; visible in project analytics.",
        "● Hard cap enforcement: the reservation system blocks any job that would exceed the cap.",
        "● Usage analytics: credit burn rate chart, top spenders by team/project, cost per video.",
        "● Stripe webhooks sync subscription status in real time — no manual plan updates needed.",
    ],
    accent_color=C_AMBER, slide_num=slide_num, total=TOTAL)

# ── Slide 12 — Multi-Channel at Scale ─────────────────────────────────────────
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
make_content_slide(sl,
    title="FOR THE ENTERPRISE OWNER  |  Multi-Channel at Scale",
    bullets=[
        "● Multiple YouTube channels per organization — each channel is independent.",
        "● Per-channel profiles stored separately:",
        "   ○ Niche profile: target audience, content category, competitor channels.",
        "   ○ Voice profile: tone, pacing, vocabulary — used by ScriptAgent for every video.",
        "   ○ Brand kit: colors, logo asset references, intro/outro preferences.",
        "● Parallel content pipelines: all channels can run AI jobs simultaneously.",
        "   ○ BullMQ concurrency: multiple workers process different channel jobs in parallel.",
        "● Channel-scoped analytics: views, watch time, CTR, and subscriber growth per channel.",
        "● Shorts Studio is channel-first: each Short is scoped to a specific channel and its brand kit.",
        "● Cross-channel reporting: BI module aggregates metrics across all org channels.",
        "● OAuth tokens encrypted at rest per channel — channel access revocable independently.",
    ],
    accent_color=C_AMBER, slide_num=slide_num, total=TOTAL)

# ── Slide 13 — Compliance at Enterprise Scale ─────────────────────────────────
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
make_content_slide(sl,
    title="FOR THE ENTERPRISE OWNER  |  Compliance at Enterprise Scale",
    bullets=[
        "● ComplianceAgent is a hard gate — no publish path bypasses it, for any team member.",
        "● ComplianceService.enforce() throws BadRequestException on failure; there is no override flag.",
        "● Compliance scoring: 0–100 scale; items scoring below 70 are blocked until revised.",
        "● ComplianceFlag types: copyright risk, advertiser-unfriendly content, misleading claims, ToS violations.",
        "● SHA-256 caching: compliance results are cached by content hash — no re-checking identical scripts.",
        "● Monetization safety score: separate dimension tracking advertiser-friendliness of every video.",
        "● FactCheckAgent citations: every factual claim in the script must link to a ResearchAgent source.",
        "   ○ Enterprise audit logs record citation provenance for every published video.",
        "● AuditLog model: every publish, approval, role change, and billing event is logged immutably.",
        "● Compliance reports available in the BI module: pass rate, common failure reasons, trend over time.",
    ],
    accent_color=C_AMBER, slide_num=slide_num, total=TOTAL)

# ══════════════════════════════════════════════════════════════════════════════
#  SECTION 4 — SUPER ADMIN / DEVELOPER OWNER PERSPECTIVE
# ══════════════════════════════════════════════════════════════════════════════

# ── Slide 14 — Super Admin Powers ─────────────────────────────────────────────
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
make_content_slide(sl,
    title="FOR THE SUPER ADMIN  |  Super Admin Powers",
    bullets=[
        "● Role = SUPER_ADMIN or OWNER → bypasses ALL plan gates in both frontend and backend.",
        "● /admin panel: view all users, manage subscriptions, override billing, reset OTPs.",
        "● Manage roles: promote/demote any user to any role (MEMBER / ORG_ADMIN / OWNER / SUPER_ADMIN).",
        "● Override billing: manually grant credits, adjust wallet balance, cancel subscriptions.",
        "● Feature flags: toggle platform-wide features via SystemConfig without a deployment.",
        "● AI ops panel: manage prompt versions, view agent token usage and cost per model.",
        "● Access to all organizations, teams, and channels — no workspace is hidden.",
        "● Railway dashboard: full backend access (NestJS logs, PostgreSQL, Redis, BullMQ queues).",
        "● Vercel dashboard: frontend deployments, env vars, preview URL management.",
        "● GitHub Actions CI: approve/retry any failed pipeline step from the Actions tab.",
        "   ○ Super Admin accounts: Digiaim_group1@iimcal.ac.in | sozialzync.vercel.app/admin",
    ],
    accent_color=C_NAVY, slide_num=slide_num, total=TOTAL)

# ── Slide 15 — System Architecture ────────────────────────────────────────────
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
make_content_slide(sl,
    title="FOR THE SUPER ADMIN  |  System Architecture",
    bullets=[
        "● Monorepo managed with pnpm workspaces + Turborepo (shared cache, parallel builds).",
        "● apps/web    — Next.js 14 App Router (frontend; deployed to Vercel).",
        "● apps/api    — NestJS (backend; deployed to Railway).",
        "● packages/agents  — AI agent implementations (stateless, idempotent).",
        "● packages/shared  — Zod schemas, shared types, utils used by both apps.",
        "● packages/prompts — Versioned prompt templates (never inline prompts in code).",
        "● packages/config  — ESLint, TypeScript, and Tailwind presets shared across packages.",
        "● infra/           — Docker Compose, IaC definitions, GitHub Actions workflows.",
        "● n8n/             — Exported workflow definitions (runtime not yet deployed).",
        "● Each NestJS module maps to one platform engine (auth, channels, projects, billing, publishing…).",
        "   ○ Turborepo caches build artifacts: only changed packages re-build in CI.",
    ],
    accent_color=C_NAVY, slide_num=slide_num, total=TOTAL)

# ── Slide 16 — Deployment Architecture ───────────────────────────────────────
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
make_content_slide(sl,
    title="FOR THE SUPER ADMIN  |  Deployment Architecture",
    bullets=[
        "● Vercel — Next.js frontend (global CDN, auto-preview URLs per push, serverless functions).",
        "   ○ URL: sozialzync.vercel.app | Auto-deploys on every push to master.",
        "   ○ Vercel env var: NEXT_PUBLIC_API_URL → points to Railway backend.",
        "● Railway — NestJS backend + PostgreSQL + Redis (persistent workers, BullMQ, WebSockets).",
        "   ○ BullMQ requires a persistent background worker — cannot run on Vercel serverless.",
        "   ○ WebSocket gateway for real-time job events also requires persistent connections.",
        "● GitHub Actions CI/CD pipeline on every push:",
        "   ○ lint → typecheck → unit tests → build → SAST (semgrep) → DAST (ZAP baseline).",
        "   ○ E2E tests (Playwright) in a separate workflow against the preview deployment.",
        "● Redis: BullMQ job queue + session cache; PostgreSQL: primary data store via Prisma ORM.",
        "● Cloudflare R2: asset storage (r2Key field on AssetVersion) — integration wiring in progress.",
    ],
    accent_color=C_NAVY, slide_num=slide_num, total=TOTAL)

# ── Slide 17 — AI Agent Ecosystem ────────────────────────────────────────────
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
make_content_slide(sl,
    title="FOR THE SUPER ADMIN  |  AI Agent Ecosystem",
    bullets=[
        "● SupervisorAgent — orchestrates the pipeline; assigns jobs, monitors retries, routes failures.",
        "● ResearchAgent   — sourced research brief with URLs, summaries, trending signals.",
        "● ScriptAgent     — hook + body + CTA script with inline fact citations.",
        "● FactCheckAgent  — verifies every claim against ResearchAgent sources; rejects hallucinations.",
        "● ComplianceAgent — scores 0–100; BLOCK flags or score < 70 halt the pipeline.",
        "● MetadataAgent   — title, description, tags, chapters, thumbnail brief.",
        "● SEOAgent        — keyword targeting, title variants, YouTube search optimization.",
        "● ThumbnailAgent  — AI thumbnail candidates via ImageAgent; user selects winner.",
        "● StoryboardAgent — scene-by-scene visual plan for video production.",
        "● AudienceAgent   — audience persona analysis, retention risk flags, engagement hooks.",
        "● CopilotAgent    — real-time chat assistant embedded in the creator workspace.",
        "   ○ All agents: stateless + idempotent. Zod-validated output. Retry → QualityControlAgent.",
    ],
    accent_color=C_NAVY, slide_num=slide_num, total=TOTAL)

# ── Slide 18 — Security & Auth ────────────────────────────────────────────────
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
make_content_slide(sl,
    title="FOR THE SUPER ADMIN  |  Security & Auth",
    bullets=[
        "● JWT access token (short-lived) + refresh token (rotation on every use, stored HttpOnly).",
        "● OTP login: email OTP via Resend API (or Gmail SMTP fallback) and SMS OTP via Twilio.",
        "   ○ OTP: 6-digit, 10-minute expiry, max 5 per window, production error if no provider configured.",
        "● Google OAuth, Apple OAuth, Facebook OAuth — all via NestJS passport strategy.",
        "● Role-based access: MEMBER / ORG_ADMIN / OWNER / SUPER_ADMIN — enforced server-side always.",
        "   ○ Frontend plan-gate is a UI affordance only; server enforces the real gate on every request.",
        "● Secrets: all keys in Railway / Vercel env vars — never committed to Git.",
        "● TOKEN_ENCRYPTION_KEY: YouTube OAuth tokens encrypted at rest in AccountLink model.",
        "● Semgrep SAST: runs on every push in GitHub Actions, catches secrets and injection patterns.",
        "● ZAP DAST baseline: runs against deployed API in CI, checks for common web vulnerabilities.",
        "● CSP headers, CORS policy, rate limiting, and Helmet middleware on all NestJS routes.",
    ],
    accent_color=C_NAVY, slide_num=slide_num, total=TOTAL)

# ── Slide 19 — Environment Variables ─────────────────────────────────────────
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
make_content_slide(sl,
    title="FOR THE SUPER ADMIN  |  Environment Variables & Configuration",
    bullets=[
        "● VERCEL (frontend only):",
        "   ○ NEXT_PUBLIC_API_URL — Railway backend URL (e.g. https://your-api.railway.app)",
        "● RAILWAY (backend — all other vars):",
        "   ○ DATABASE_URL       — PostgreSQL connection string (Railway plugin).",
        "   ○ REDIS_URL          — Redis connection string (Railway plugin).",
        "   ○ JWT_SECRET         — random 256-bit string for signing access tokens.",
        "   ○ JWT_REFRESH_SECRET — separate secret for refresh tokens.",
        "   ○ TOKEN_ENCRYPTION_KEY — AES key for encrypting stored OAuth tokens.",
        "   ○ RESEND_API_KEY + RESEND_FROM — email OTP delivery via Resend.",
        "   ○ TWILIO_ACCOUNT_SID + TWILIO_AUTH_TOKEN + TWILIO_FROM — SMS OTP.",
        "   ○ ANTHROPIC_API_KEY / OPENAI_API_KEY / GOOGLE_AI_API_KEY — AI providers.",
        "   ○ STRIPE_SECRET_KEY + STRIPE_WEBHOOK_SECRET — billing.",
        "   ○ GOOGLE_CLIENT_ID + GOOGLE_CLIENT_SECRET — YouTube OAuth.",
    ],
    accent_color=C_NAVY, slide_num=slide_num, total=TOTAL)

# ══════════════════════════════════════════════════════════════════════════════
#  SECTION 5 — DEVELOPER PERSPECTIVE
# ══════════════════════════════════════════════════════════════════════════════

# ── Slide 20 — Tech Stack Overview ───────────────────────────────────────────
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
make_content_slide(sl,
    title="FOR THE DEVELOPER  |  Tech Stack Overview",
    bullets=[
        "● FRONTEND (apps/web — Vercel)",
        "   ○ Next.js 14 App Router, TypeScript strict, Tailwind CSS, TanStack Query v5.",
        "   ○ Lucide React icons, Radix UI primitives, Framer Motion animations.",
        "   ○ Mobile-first: bottom nav on mobile, sidebar drawer, responsive breakpoints.",
        "● BACKEND (apps/api — Railway)",
        "   ○ NestJS (modular), Prisma ORM, PostgreSQL, Redis, BullMQ async jobs.",
        "   ○ Passport.js (JWT + OAuth strategies), Helmet, class-validator, class-transformer.",
        "   ○ WebSocket gateway (Socket.io) for real-time job progress events.",
        "● AI PROVIDERS",
        "   ○ Claude (Anthropic) — primary model for scripting and compliance.",
        "   ○ OpenAI GPT — fallback provider via shared aiClient.",
        "   ○ Google Gemini — fallback provider; also used for YouTube channel analytics.",
        "● TOOLING: pnpm workspaces, Turborepo, GitHub Actions, Playwright E2E, Jest unit tests.",
    ],
    accent_color=C_TEAL, slide_num=slide_num, total=TOTAL)

# ── Slide 21 — Monorepo Structure ────────────────────────────────────────────
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
make_content_slide(sl,
    title="FOR THE DEVELOPER  |  Monorepo Structure",
    bullets=[
        "● creatorforce-ai/           ← root (pnpm workspaces + Turborepo)",
        "   ○ apps/web/               ← Next.js 14 frontend (Vercel)",
        "   ○ apps/api/               ← NestJS backend (Railway)",
        "   ○ packages/agents/        ← AI agent implementations",
        "   ○ packages/shared/        ← Zod schemas, types, utils (shared by web + api)",
        "   ○ packages/prompts/       ← versioned prompt templates (never inline in code)",
        "   ○ packages/config/        ← ESLint, tsconfig, Tailwind presets",
        "   ○ infra/                  ← Docker Compose, GitHub Actions workflows",
        "   ○ n8n/                    ← exported n8n workflow definitions",
        "   ○ docs/                   ← all design docs (read before coding)",
        "   ○ CLAUDE.md               ← AI coding agent operating contract",
        "● Turborepo pipeline: build depends on ^build; test depends on ^build.",
        "● pnpm workspace protocol: packages reference each other via workspace:*.",
    ],
    accent_color=C_TEAL, slide_num=slide_num, total=TOTAL)

# ── Slide 22 — Frontend Architecture ─────────────────────────────────────────
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
make_content_slide(sl,
    title="FOR THE DEVELOPER  |  Frontend Architecture",
    bullets=[
        "● Next.js App Router — Server Components by default; Client Components only when interactive.",
        "● Route groups: (dash) for authenticated dashboard, (auth) for login/signup flows.",
        "● /api/proxy route: frontend API calls go through Next.js proxy → Railway backend (avoids CORS).",
        "● TanStack Query: all data fetching; shared queryKey conventions; 2-minute staleTime.",
        "● State management: React useState + useEffect; no global store (context for auth only).",
        "● Hydration safety: localStorage reads always inside useEffect, never at render time.",
        "   ○ Pattern: useState(defaultValue) → useEffect(() => { setValue(localStorage.get()) }, []).",
        "● Mobile responsive: bottom nav bar on mobile, sidebar drawer (Sheet), touch targets ≥ 44px.",
        "● Plan gating: PlanGate component + usePlanGate hook; reads JWT from localStorage after mount.",
        "● Styling: Tailwind CSS utility classes; design tokens via Tailwind config (brand purple = #7C3AED).",
        "● Type safety: TypeScript strict; Zod schemas from packages/shared used for API response parsing.",
    ],
    accent_color=C_TEAL, slide_num=slide_num, total=TOTAL)

# ── Slide 23 — Backend Architecture ──────────────────────────────────────────
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
make_content_slide(sl,
    title="FOR THE DEVELOPER  |  Backend Architecture",
    bullets=[
        "● NestJS modules — one module per engine: auth, channels, projects, jobs, wallet, billing, orgs, publishing, analytics.",
        "● Async jobs: anything > 2s or calling an external AI/video/music provider runs as a BullMQ job.",
        "   ○ AGENT_QUEUE: ResearchJob, ScriptJob, FactCheckJob, ComplianceJob, MetadataJob, SEOJob, PublishJob.",
        "   ○ SHORTS_QUEUE: ShortsAnalyzeJob and dedicated Shorts Studio job types.",
        "● WebSocket gateway (Socket.io): emits job:progress, job:complete, job:error events to the frontend.",
        "● Guards: JwtAuthGuard (all routes), DeveloperKeyGuard (external API), RolesGuard (admin).",
        "● Middleware: Helmet (security headers), CORS, rate limiter, request logger, Sentry error capture.",
        "● Error handling: typed domain errors thrown, never swallowed; all surface to Sentry.",
        "● Observability: Prometheus metrics via prom-client on every route; AuditLog for compliance trails.",
        "● Prisma ORM: type-safe DB access; migrations via prisma migrate deploy in CI.",
        "● Zod validation: at every boundary — API input (class-validator), agent output, env vars (zod.parse).",
    ],
    accent_color=C_TEAL, slide_num=slide_num, total=TOTAL)

# ── Slide 24 — Agent Development Rules ───────────────────────────────────────
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
make_content_slide(sl,
    title="FOR THE DEVELOPER  |  Agent Development Rules",
    bullets=[
        "● Step 1: Define input + output Zod schemas in packages/shared — never inline in agent code.",
        "● Step 2: Pull the prompt from packages/prompts (versioned file) — never hardcode a large prompt.",
        "● Step 3: Call the AI model via the shared aiClient (handles retries, fallback provider, token accounting, tracing).",
        "● Step 4: Validate output against the output Zod schema.",
        "   ○ On schema failure: retry up to MAX_AGENT_RETRIES, then route to QualityControlAgent.",
        "● Step 5: Emit a structured trace event (agent name, model, tokens, latency, cost estimate).",
        "● Step 6: Agents must remain stateless and idempotent — SupervisorAgent orchestrates sequencing.",
        "● Provider fallback order: Claude (Anthropic) → OpenAI → Gemini (configured in aiClient).",
        "● Token accounting: all provider calls log tokens + estimated cost to AgentLog model.",
        "● ComplianceAgent is special: its output is also stored in ComplianceResult and is a hard gate.",
        "   ○ ComplianceService.enforce() throws BadRequestException — the pipeline does NOT continue on failure.",
    ],
    accent_color=C_TEAL, slide_num=slide_num, total=TOTAL)

# ── Slide 25 — Database Schema ────────────────────────────────────────────────
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
make_content_slide(sl,
    title="FOR THE DEVELOPER  |  Database Schema Highlights",
    bullets=[
        "● User, AccountLink (OAuth tokens encrypted), Session, OtpCode — auth models.",
        "● Subscription (FREE/STARTER/PRO/AGENCY), Wallet (user or org), CreditLedger, CreditLot — billing.",
        "   ○ CreditReservation: hold-and-settle pattern; BudgetPeriod: monthly cap per org/team.",
        "● Organization, OrgMembership, Team, TeamMembership — org/team hierarchy.",
        "● Channel, LibraryVideo, LibraryPlaylist, AccountLink — YouTube channel models.",
        "● Project, Script, Approval, ComplianceResult, ComplianceFlag — content pipeline.",
        "● Job (BullMQ job record), AgentLog (token usage, cost, latency per agent call).",
        "● ImportedVideo, TranscriptSegment, VideoScene, TopicSegment, Chapter — Shorts Studio.",
        "● ShortClip, ShortsTimeline, SocialContent — Shorts generation models.",
        "● Asset, AssetVersion (r2Key, provider field) — media pipeline with R2 storage keys.",
        "● ReferralCode, TrialGrant, DeveloperKey, DeveloperWebhook, AuditLog, SystemConfig — growth/ops.",
    ],
    accent_color=C_TEAL, slide_num=slide_num, total=TOTAL)

# ── Slide 26 — CI/CD & Build Pipeline ────────────────────────────────────────
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
make_content_slide(sl,
    title="FOR THE DEVELOPER  |  CI/CD & Build Pipeline",
    bullets=[
        "● GitHub Actions — triggers on every push to master and on pull requests.",
        "● Pipeline steps (in order):",
        "   ○ lint         — ESLint across all packages (Turborepo cached).",
        "   ○ typecheck    — tsc --noEmit strict mode across all packages.",
        "   ○ unit tests   — Jest; co-located *.spec.ts files; coverage threshold enforced.",
        "   ○ build        — Turborepo build (only changed packages re-build).",
        "   ○ SAST         — Semgrep rules (secrets, injection, security anti-patterns).",
        "   ○ DAST         — OWASP ZAP baseline scan against deployed preview API.",
        "   ○ E2E          — Playwright tests against Vercel preview deployment.",
        "● Vercel: auto-deploys on every push; preview URL per PR; production on master merge.",
        "● Railway: auto-deploys on env var change or when triggered by GitHub Actions.",
        "● prisma migrate deploy runs as part of the Railway deploy hook (never manual).",
        "● Branch protection: PR must pass all checks before merge; no force-push to master.",
    ],
    accent_color=C_TEAL, slide_num=slide_num, total=TOTAL)

# ── Slide 27 — Coding Conventions ────────────────────────────────────────────
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
make_content_slide(sl,
    title="FOR THE DEVELOPER  |  Coding Conventions",
    bullets=[
        "● TypeScript strict mode everywhere — no any without a // @reason: comment.",
        "● Naming: PascalCase types/classes | camelCase vars/functions | SCREAMING_SNAKE env vars | kebab-case files.",
        "● Zod at every boundary: API input, agent output, environment variables (zod.parse on startup).",
        "● BullMQ for async: any operation > 2s or calling an external provider must be a job, not inline.",
        "● Errors: throw typed domain errors; never swallow exceptions; all errors surface to Sentry.",
        "● Tests: co-locate *.spec.ts with source. New features require unit tests; pipelines need integration tests.",
        "● Comments: only when the WHY is non-obvious. No docblocks. No task references ('fix for issue #123').",
        "● Conventional Commits: feat: / fix: / chore: / docs: / refactor: / test: — one logical change per PR.",
        "● No inline prompts: large prompts live in packages/prompts with version numbers.",
        "● No feature flags or backwards-compat shims unless explicitly required — just change the code.",
        "● Security: no command injection, no XSS, no SQL injection — Prisma parameterized queries always.",
    ],
    accent_color=C_TEAL, slide_num=slide_num, total=TOTAL)

# ══════════════════════════════════════════════════════════════════════════════
#  SECTION 6 — CLOSING
# ══════════════════════════════════════════════════════════════════════════════

# ── Slide 28 — Roadmap ────────────────────────────────────────────────────────
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
make_content_slide(sl,
    title="ROADMAP  |  What's Next",
    bullets=[
        "● Phase 1 — Core Platform (COMPLETED)",
        "   ○ Auth, channel connect, content pipeline, compliance gate, billing, developer portal, CI.",
        "● Phase 2 — Shorts Studio (COMPLETED)",
        "   ○ Channel-first flow, transcript analysis, clip recommendations, timeline editor, social factory.",
        "● Phase 3 — Enterprise & Growth (COMPLETED)",
        "   ○ Org/team workspaces, referral program, trial engine, BI analytics module.",
        "● Phase 4 — Media Pipeline (IN PROGRESS)",
        "   ○ Built: Voice, Image, Music, Video, Subtitle agents; ffmpeg render pipeline; Asset/AssetVersion models.",
        "   ○ Still needed: Cloudflare R2 SDK wiring, external video/music provider integrations (Veo/Kling/Suno).",
        "● Phase 5 — Scale & Ops (PLANNED)",
        "   ○ n8n workflow automation, multi-region deployment, i18n beyond English, accessibility audit.",
        "   ○ Stripe production keys, horizontal BullMQ worker scaling, Postgres read replicas.",
    ],
    accent_color=C_PURPLE2, slide_num=slide_num, total=TOTAL)

# ── Slide 29 — Key Links & Contacts ──────────────────────────────────────────
slide_num += 1
sl = prs.slides.add_slide(blank_layout)
make_content_slide(sl,
    title="KEY LINKS & CONTACTS",
    bullets=[
        "● App URL:         https://sozialzync.vercel.app",
        "● Admin Panel:     https://sozialzync.vercel.app/admin",
        "● GitHub Repo:     https://github.com/MoshayM/AI-Creatorforce",
        "● Super Admin:     Digiaim_group1@iimcal.ac.in",
        "● OTP Setup — configure these Railway env vars for OTP delivery to work:",
        "   ○ Email (Resend): RESEND_API_KEY + RESEND_FROM",
        "   ○ Email (SMTP):   SMTP_HOST + SMTP_PORT + SMTP_USER + SMTP_PASS + SMTP_FROM",
        "   ○ SMS (Twilio):   TWILIO_ACCOUNT_SID + TWILIO_AUTH_TOKEN + TWILIO_FROM",
        "● For a full setup guide, see: OTP-SETUP.txt in the repo root.",
        "● Frontend deployed on Vercel — auto-redeploys on every push to master.",
        "● Backend deployed on Railway — redeploys when env vars change (~30 seconds).",
        "   ○ Full architecture details: docs/architecture.md | docs/deployment.md",
    ],
    accent_color=C_PURPLE2, slide_num=slide_num, total=TOTAL)

# ── Save ───────────────────────────────────────────────────────────────────────
out_path = r"D:\project\creatorforce-ai\AI-CreatorForce-Presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
