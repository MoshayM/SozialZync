"""
make_superadmin_pptx.py — AI CreatorForce Super Admin Guide
Generates CreatorForce-SuperAdmin-Guide.pptx (24 slides)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy
from lxml import etree

# ---------------------------------------------------------------------------
# Color constants
# ---------------------------------------------------------------------------
NAVY   = RGBColor(0x1E, 0x1B, 0x4B)   # #1E1B4B  header bar
GOLD   = RGBColor(0xF5, 0x9E, 0x0B)   # #F59E0B  accent
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)   # #FFFFFF
LIGHT  = RGBColor(0xC7, 0xD2, 0xFE)   # #C7D2FE  light accent
BG     = RGBColor(0x0F, 0x17, 0x2A)   # #0F172A  dark bg
SLATE  = RGBColor(0x94, 0xA3, 0xB8)   # #94A3B8  muted text
ROW_A  = RGBColor(0x1E, 0x29, 0x3B)   # #1E293B  table row alt
ROW_B  = RGBColor(0x0F, 0x17, 0x2A)   # #0F172A  table row base

# ---------------------------------------------------------------------------
# Slide dimensions
# ---------------------------------------------------------------------------
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]

# ---------------------------------------------------------------------------
# Helper utilities
# ---------------------------------------------------------------------------

def add_rect(slide, left, top, width, height, fill_color, line=False):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if not line:
        shape.line.width = 0
    return shape


def add_textbox(slide, left, top, width, height,
                text, font_size, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def add_textbox_lines(slide, left, top, width, height,
                      lines, font_size, color=WHITE,
                      align=PP_ALIGN.LEFT, bold=False, line_spacing=None):
    """Add a textbox with multiple paragraphs (one per item in lines list)."""
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for line_text in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if line_spacing is not None:
            p.line_spacing = Pt(line_spacing)
        run = p.add_run()
        run.text = line_text
        run.font.size  = Pt(font_size)
        run.font.bold  = bold
        run.font.color.rgb = color
    return txb


def add_slide_number(slide, slide_num):
    """Add slide number at bottom right in gold."""
    add_textbox(slide,
                Inches(12.3), Inches(7.0), Inches(0.9), Inches(0.4),
                str(slide_num), 11, bold=False, color=GOLD, align=PP_ALIGN.RIGHT)


def add_header_bar(slide, title_text):
    """Standard header: navy bar + gold accent line + white title text."""
    # Navy bar
    add_rect(slide, 0, 0, W, Inches(1.2), NAVY)
    # Gold accent bar
    add_rect(slide, 0, Inches(1.2), W, Inches(0.07), GOLD)
    # Title in header
    add_textbox(slide,
                Inches(0.3), Inches(0.1),
                Inches(12.7), Inches(1.0),
                title_text, 28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)


def fill_slide_bg(slide, color=BG):
    """Fill entire slide background."""
    add_rect(slide, 0, 0, W, H, color)


def add_section_break(slide, title, subtitle=None):
    """Full dark bg section break slide."""
    fill_slide_bg(slide)
    # Gold accent line top
    add_rect(slide, 0, 0, W, Inches(0.08), GOLD)
    # Large title centered
    add_textbox(slide,
                Inches(0.5), Inches(2.4),
                Inches(12.33), Inches(1.5),
                title, 44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        add_textbox(slide,
                    Inches(0.5), Inches(4.2),
                    Inches(12.33), Inches(0.7),
                    subtitle, 20, bold=False, color=GOLD, align=PP_ALIGN.CENTER)


def set_cell_bg(cell, rgb: RGBColor):
    """Set table cell background color via lxml."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = etree.SubElement(tcPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    srgbClr   = etree.SubElement(solidFill, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
    srgbClr.set('val', '{:02X}{:02X}{:02X}'.format(rgb[0], rgb[1], rgb[2]))


def set_cell_text(cell, text, font_size=13, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    tf = cell.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color


# ===========================================================================
# SLIDE 1 — Title slide
# ===========================================================================
sl = prs.slides.add_slide(BLANK)
fill_slide_bg(sl)

# Gold top accent bar
add_rect(sl, 0, 0, W, Emu(91440), GOLD)   # 91440 EMU ≈ 8px visual at 96dpi

# Main title
add_textbox(sl,
            Inches(0.5), Inches(2.0),
            Inches(12.33), Inches(0.9),
            "AI CreatorForce",
            52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Subtitle
add_textbox(sl,
            Inches(0.5), Inches(3.0),
            Inches(12.33), Inches(0.6),
            "Super Admin & Platform Owner Reference",
            24, bold=False, color=LIGHT, align=PP_ALIGN.CENTER)

# Crown / badge line
add_textbox(sl,
            Inches(0.5), Inches(3.8),
            Inches(12.33), Inches(0.45),
            "★  FULL PLATFORM ACCESS  ★",
            16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# URL line
add_textbox(sl,
            Inches(0.5), Inches(5.0),
            Inches(12.33), Inches(0.45),
            "sozialzync.vercel.app/admin  |  github.com/MoshayM/AI-Creatorforce",
            13, bold=False, color=SLATE, align=PP_ALIGN.CENTER)

add_slide_number(sl, 1)

# ===========================================================================
# SLIDE 2 — Section break: YOUR ROLE AS SUPER ADMIN
# ===========================================================================
sl = prs.slides.add_slide(BLANK)
add_section_break(sl,
                  "YOUR ROLE AS SUPER ADMIN",
                  "Owner · Developer · Platform Operator")
add_slide_number(sl, 2)

# ===========================================================================
# SLIDE 3 — What Super Admin Access Means
# ===========================================================================
sl = prs.slides.add_slide(BLANK)
fill_slide_bg(sl)
add_header_bar(sl, "What Super Admin Access Means")

bullets = [
    "● You are the platform owner — you built and deployed this system",
    "● JWT role = SUPER_ADMIN or OWNER → bypasses ALL plan gates",
    "● Full access to every feature regardless of plan tier",
    "● Access to /admin panel — only visible to SUPER_ADMIN/OWNER",
    "● Can view, manage, and impersonate any user account",
    "● Can override billing, credit balances, and plan assignments",
    "● Responsible for uptime, security, and compliance of the platform",
]
add_textbox_lines(sl,
                  Inches(0.4), Inches(1.45),
                  Inches(12.5), Inches(5.7),
                  bullets, 15, color=WHITE, line_spacing=28)
add_slide_number(sl, 3)

# ===========================================================================
# SLIDE 4 — Super Admin vs Other Roles (table)
# ===========================================================================
sl = prs.slides.add_slide(BLANK)
fill_slide_bg(sl)
add_header_bar(sl, "Super Admin vs Other Roles")

table_data = [
    ["Role",        "Plan Gate",       "Admin Panel",   "Billing Override"],
    ["MEMBER",      "Full gates",      "No",            "No"],
    ["ORG_ADMIN",   "Org features",    "No",            "Org wallet only"],
    ["OWNER",       "Bypass all",      "Yes (org)",     "Yes (org)"],
    ["SUPER_ADMIN", "Bypass ALL",      "Yes (global)",  "Yes (global)"],
]

rows, cols = len(table_data), len(table_data[0])
tbl = sl.shapes.add_table(rows, cols,
                           Inches(0.4), Inches(1.5),
                           Inches(12.5), Inches(4.5)).table

col_widths = [Inches(2.8), Inches(3.2), Inches(3.2), Inches(3.3)]
for i, w in enumerate(col_widths):
    tbl.columns[i].width = w

for r_idx, row in enumerate(table_data):
    for c_idx, cell_text in enumerate(row):
        cell = tbl.cell(r_idx, c_idx)
        if r_idx == 0:
            set_cell_bg(cell, NAVY)
            set_cell_text(cell, cell_text, font_size=14, bold=True,
                          color=WHITE, align=PP_ALIGN.CENTER)
        else:
            bg = ROW_A if r_idx % 2 == 1 else ROW_B
            set_cell_bg(cell, bg)
            # Highlight SUPER_ADMIN row in gold
            txt_color = GOLD if r_idx == 4 else WHITE
            set_cell_text(cell, cell_text, font_size=13, bold=(r_idx == 4),
                          color=txt_color, align=PP_ALIGN.CENTER)

add_slide_number(sl, 4)

# ===========================================================================
# SLIDE 5 — The /admin Panel
# ===========================================================================
sl = prs.slides.add_slide(BLANK)
fill_slide_bg(sl)
add_header_bar(sl, "The /admin Panel")

bullets = [
    "● Access: sozialzync.vercel.app/admin  (role-gated — SUPER_ADMIN/OWNER only)",
    "● User Management: list all users, search, view profiles, assign roles",
    "● Billing Control: view wallet balances, grant/revoke credits, adjust plans",
    "● Organization Overview: all orgs, member counts, billing status",
    "● Job Queue Monitor: BullMQ job states, failed jobs, retry controls",
    "● Analytics: platform-wide usage, top users, revenue metrics",
    "● Audit Log: every significant action logged with user + timestamp",
]
add_textbox_lines(sl,
                  Inches(0.4), Inches(1.45),
                  Inches(12.5), Inches(5.7),
                  bullets, 15, color=WHITE, line_spacing=28)
add_slide_number(sl, 5)

# ===========================================================================
# SLIDE 6 — Section break: PLATFORM DEPLOYMENT
# ===========================================================================
sl = prs.slides.add_slide(BLANK)
add_section_break(sl, "PLATFORM DEPLOYMENT")
add_slide_number(sl, 6)

# ===========================================================================
# SLIDE 7 — Infrastructure Overview (two-column)
# ===========================================================================
sl = prs.slides.add_slide(BLANK)
fill_slide_bg(sl)
add_header_bar(sl, "Infrastructure Overview")

# Left column — Frontend (Vercel)
left_lines = [
    "FRONTEND (Vercel)",
    "",
    "● URL: sozialzync.vercel.app",
    "● Auto-deploys on git push to master",
    "● Preview URLs for every PR branch",
    "● Only one env var needed: NEXT_PUBLIC_API_URL",
    "● Free Vercel plan sufficient for most usage",
]
left_tb = sl.shapes.add_textbox(Inches(0.4), Inches(1.45), Inches(5.8), Inches(5.5))
left_tf = left_tb.text_frame
left_tf.word_wrap = True
for i, line in enumerate(left_lines):
    p = left_tf.paragraphs[0] if i == 0 else left_tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = line
    run.font.size = Pt(14 if i == 0 else 14)
    run.font.bold = (i == 0)
    run.font.color.rgb = GOLD if i == 0 else WHITE

# Right column — Backend (Railway)
right_lines = [
    "BACKEND (Railway)",
    "",
    "● NestJS API server (persistent)",
    "● BullMQ worker process (persistent)",
    "● PostgreSQL database (Railway plugin)",
    "● Redis for queues/cache (Railway plugin)",
    "● All env vars configured in Railway Variables tab",
]
right_tb = sl.shapes.add_textbox(Inches(6.8), Inches(1.45), Inches(6.1), Inches(5.5))
right_tf = right_tb.text_frame
right_tf.word_wrap = True
for i, line in enumerate(right_lines):
    p = right_tf.paragraphs[0] if i == 0 else right_tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = line
    run.font.size = Pt(14)
    run.font.bold = (i == 0)
    run.font.color.rgb = GOLD if i == 0 else WHITE

# Vertical divider
add_rect(sl, Inches(6.55), Inches(1.5), Inches(0.03), Inches(5.5), NAVY)

add_slide_number(sl, 7)

# ===========================================================================
# SLIDE 8 — Required Environment Variables — Railway (Core)
# ===========================================================================
sl = prs.slides.add_slide(BLANK)
fill_slide_bg(sl)
add_header_bar(sl, "Required Environment Variables — Railway")

lines = [
    "Core variables (MUST SET — backend will not start without these):",
    "",
    "● DATABASE_URL — Railway auto-sets when PostgreSQL plugin is added",
    "● REDIS_URL — Railway auto-sets when Redis plugin is added",
    "● JWT_SECRET — random 64-char string  (generate: openssl rand -hex 32)",
    "● JWT_REFRESH_SECRET — separate random 64-char string",
    "● TOKEN_ENCRYPTION_KEY — min 32 chars (JWE key for OAuth token encryption)",
    "● SUPER_ADMIN_EMAILS — comma-separated list of super-admin email addresses",
    "● OWNER_EMAILS — comma-separated list of owner email addresses",
    "● NODE_ENV=production",
    "● CORS_ORIGIN=https://sozialzync.vercel.app",
]
add_textbox_lines(sl,
                  Inches(0.4), Inches(1.45),
                  Inches(12.5), Inches(5.7),
                  lines, 14, color=WHITE, line_spacing=26)
add_slide_number(sl, 8)

# ===========================================================================
# SLIDE 9 — Required Environment Variables — AI Providers
# ===========================================================================
sl = prs.slides.add_slide(BLANK)
fill_slide_bg(sl)
add_header_bar(sl, "Required Environment Variables — AI Providers")

lines = [
    "AI API Keys (configure at least one — Claude is primary):",
    "  ● ANTHROPIC_API_KEY — Claude (primary AI provider)",
    "  ● OPENAI_API_KEY — GPT-4 (fallback provider)",
    "  ● GEMINI_API_KEY — Gemini (second fallback)",
    "",
    "OTP Delivery (configure at least one — Resend recommended):",
    "  ● RESEND_API_KEY + RESEND_FROM — email OTP (free, 3,000 emails/month)",
    "  ● SMTP_HOST + SMTP_PORT + SMTP_USER + SMTP_PASS + SMTP_FROM — SMTP fallback",
    "  ● TWILIO_ACCOUNT_SID + TWILIO_AUTH_TOKEN + TWILIO_FROM — SMS OTP",
    "",
    "Payments:",
    "  ● STRIPE_SECRET_KEY + STRIPE_WEBHOOK_SECRET",
    "",
    "Observability (optional but recommended):",
    "  ● SENTRY_DSN — error tracking for both API and frontend",
]
add_textbox_lines(sl,
                  Inches(0.4), Inches(1.45),
                  Inches(12.5), Inches(5.7),
                  lines, 13.5, color=WHITE, line_spacing=25)
add_slide_number(sl, 9)

# ===========================================================================
# SLIDE 10 — Deployment Procedure Step by Step
# ===========================================================================
sl = prs.slides.add_slide(BLANK)
fill_slide_bg(sl)
add_header_bar(sl, "Deployment Procedure — Step by Step")

lines = [
    "● 1. Push code to master branch → GitHub Actions CI runs automatically",
    "● 2. CI pipeline: typecheck → lint → unit tests → build (all must pass)",
    "● 3. CI also runs: SAST (Semgrep), dependency audit, ZAP baseline, E2E tests",
    "● 4. On CI success: Vercel webhook triggers frontend deploy  (~2 minutes)",
    "● 5. For backend: Railway auto-redeploys on env var change or manual trigger",
    "       Run: prisma migrate deploy  before restarting API for schema changes",
    "● 6. Verify: check sozialzync.vercel.app loads, login works, API responds at /health",
    "● 7. Monitor: Railway logs for backend errors, Vercel logs for frontend issues",
    "● 8. Rollback: Vercel → 'Promote' previous deployment; Railway: revert Railway deploy",
]
add_textbox_lines(sl,
                  Inches(0.4), Inches(1.45),
                  Inches(12.5), Inches(5.7),
                  lines, 14.5, color=WHITE, line_spacing=27)
add_slide_number(sl, 10)

# ===========================================================================
# SLIDE 11 — Section break: SECURITY & COMPLIANCE
# ===========================================================================
sl = prs.slides.add_slide(BLANK)
add_section_break(sl, "SECURITY & COMPLIANCE")
add_slide_number(sl, 11)

# ===========================================================================
# SLIDE 12 — Security Architecture
# ===========================================================================
sl = prs.slides.add_slide(BLANK)
fill_slide_bg(sl)
add_header_bar(sl, "Security Architecture")

lines = [
    "● JWT: 15-minute access tokens, 7-day refresh tokens (rotating pattern)",
    "● Refresh tokens: stored in HTTP-only SameSite cookie, not localStorage",
    "● OTP codes: hashed in DB, expire in 10 minutes, single-use only",
    "● Passwords: bcrypt hashed with cost factor — never stored in plaintext",
    "● OAuth tokens: AES-256-GCM encrypted at rest using TOKEN_ENCRYPTION_KEY",
    "● API keys: env vars only — never in code, never in git history",
    "● CORS: restricted to CORS_ORIGIN env var value (sozialzync.vercel.app)",
    "● HTTP headers: CSP, HSTS, X-Frame-Options DENY, nosniff (Helmet 8)",
    "● Rate limiting: Redis-backed on OTP, login, register, refresh endpoints",
    "● SAST: Semgrep custom rules run in CI on every push — ERROR level blocks merge",
    "● DAST: OWASP ZAP baseline scan on every push — HIGH findings block merge",
    "● Dependency audit: pnpm audit --audit-level=high blocks on HIGH/CRITICAL CVEs",
]
add_textbox_lines(sl,
                  Inches(0.4), Inches(1.45),
                  Inches(12.5), Inches(5.7),
                  lines, 13.5, color=WHITE, line_spacing=25)
add_slide_number(sl, 12)

# ===========================================================================
# SLIDE 13 — The Compliance Hard Gate
# ===========================================================================
sl = prs.slides.add_slide(BLANK)
fill_slide_bg(sl)
add_header_bar(sl, "The Compliance Hard Gate")

lines = [
    "● ComplianceAgent runs on EVERY content item before it can be published",
    "● No code path exists that skips ComplianceAgent — enforced in CI by Semgrep",
    "● Checks: advertiser-friendliness, copyright risk, YouTube community guidelines",
    "● Assigns Monetization Safety Score (0–100) to each content item",
    "● Score below threshold → content is blocked; user notified with specific reasons",
    "● All compliance decisions are written to the AuditLog table (append-only)",
    "● Publishing endpoint rejects any request without a valid passed ComplianceResult",
    "● ComplianceResult is re-verified by the publishing worker (double-check gate)",
    "● Editing a script after approval resets both compliance and human-approval gates",
]
add_textbox_lines(sl,
                  Inches(0.4), Inches(1.45),
                  Inches(12.5), Inches(5.7),
                  lines, 15, color=WHITE, line_spacing=28)
add_slide_number(sl, 13)

# ===========================================================================
# SLIDE 14 — Human-in-the-Loop Enforcement
# ===========================================================================
sl = prs.slides.add_slide(BLANK)
fill_slide_bg(sl)
add_header_bar(sl, "Human-in-the-Loop Enforcement")

lines = [
    "● The API's /publishing/publish endpoint requires explicit human approval",
    "● Approval is a separate API call — publishing cannot happen in a single step",
    "● Scheduled auto-publish: ONLY if user explicitly enabled it AND item passed compliance",
    "● Approval has an expiry — expired approval blocks publish; new approval required",
    "● Super Admin can view the publish queue and cancel any pending publishes",
    "● Every publish event is logged: user, channel, video ID, timestamp, compliance score",
    "● YouTube Data API upload only happens after: compliance pass + human approval",
    "● Approval status in DB: PENDING / APPROVED / REJECTED / EXPIRED",
]
add_textbox_lines(sl,
                  Inches(0.4), Inches(1.45),
                  Inches(12.5), Inches(5.7),
                  lines, 15, color=WHITE, line_spacing=28)
add_slide_number(sl, 14)

# ===========================================================================
# SLIDE 15 — Section break: USER & BILLING MANAGEMENT
# ===========================================================================
sl = prs.slides.add_slide(BLANK)
add_section_break(sl, "USER & BILLING MANAGEMENT")
add_slide_number(sl, 15)

# ===========================================================================
# SLIDE 16 — Managing Users
# ===========================================================================
sl = prs.slides.add_slide(BLANK)
fill_slide_bg(sl)
add_header_bar(sl, "Managing Users")

lines = [
    "● View all users in /admin panel — searchable, filterable by role and plan",
    "● Change user role: MEMBER ↔ ORG_ADMIN ↔ OWNER ↔ SUPER_ADMIN",
    "● Reset password / resend OTP verification: directly from admin panel",
    "● Suspend/reactivate accounts: instantly blocks all API requests for that user",
    "● View user's projects, active jobs, and published videos",
    "● Impersonate (view-only mode): see exactly what the user sees without modifying data",
    "● Delete user: GDPR-compliant — scrubs PII, anonymizes associated audit log rows",
    "● Fraud control: set Wallet.rechargesFrozen=true to block wallet top-ups",
]
add_textbox_lines(sl,
                  Inches(0.4), Inches(1.45),
                  Inches(12.5), Inches(5.7),
                  lines, 15, color=WHITE, line_spacing=27)
add_slide_number(sl, 16)

# ===========================================================================
# SLIDE 17 — Billing & Credit Control
# ===========================================================================
sl = prs.slides.add_slide(BLANK)
fill_slide_bg(sl)
add_header_bar(sl, "Billing & Credit Control")

lines = [
    "● Grant credits to any wallet: select user → Wallet → Grant Credits",
    "● Credit types: Trial / Purchased / Bonus / Referral — each tracked separately",
    "● Set expiry date for granted credits (optional — null means never expires)",
    "● Override plan tier: assign Enterprise to any user manually from admin panel",
    "● View all Stripe subscriptions and full payment history per user",
    "● Process refunds: via Stripe Dashboard (not from admin panel directly)",
    "● Credit ledger: every debit/credit transaction logged with entry type and reason",
    "● CreditLedger is append-only — corrections are new ADJUSTMENT rows, never edits",
    "● Hard budget cap: per-user monthly spend cap enforced fail-closed in WalletService",
]
add_textbox_lines(sl,
                  Inches(0.4), Inches(1.45),
                  Inches(12.5), Inches(5.7),
                  lines, 14.5, color=WHITE, line_spacing=26)
add_slide_number(sl, 17)

# ===========================================================================
# SLIDE 18 — Organization Management
# ===========================================================================
sl = prs.slides.add_slide(BLANK)
fill_slide_bg(sl)
add_header_bar(sl, "Organization Management")

lines = [
    "● View all organizations (Enterprise plan users) in the admin panel",
    "● See: org name, owner, member count, wallet balance, credit usage",
    "● Adjust org wallet: grant or revoke org-level credits from admin panel",
    "● Add/remove org members directly from admin panel",
    "● Set monthly budget cap for org with hardCap=true (prevents runaway spend)",
    "● View org's channels, projects, and publish history",
    "● Organization statuses: ACTIVE / SUSPENDED — toggle from admin panel",
    "● Dissolve org: reassigns members to personal accounts",
    "● BudgetPeriod table: per-org, per-period credit allocation and hard-cap enforcement",
]
add_textbox_lines(sl,
                  Inches(0.4), Inches(1.45),
                  Inches(12.5), Inches(5.7),
                  lines, 14.5, color=WHITE, line_spacing=26)
add_slide_number(sl, 18)

# ===========================================================================
# SLIDE 19 — Section break: MONITORING & OPERATIONS
# ===========================================================================
sl = prs.slides.add_slide(BLANK)
add_section_break(sl, "MONITORING & OPERATIONS")
add_slide_number(sl, 19)

# ===========================================================================
# SLIDE 20 — Queue & Job Monitoring
# ===========================================================================
sl = prs.slides.add_slide(BLANK)
fill_slide_bg(sl)
add_header_bar(sl, "Queue & Job Monitoring")

lines = [
    "● BullMQ Board: visual dashboard for all job queues (accessible from admin panel)",
    "● Queues: agent-jobs, publish-jobs, import-jobs, analytics-jobs, automation-tick",
    "● View per-queue counts: waiting / active / completed / failed / delayed",
    "● Retry failed jobs manually from the dashboard — no code deploy needed",
    "● Pause/resume queues during maintenance windows",
    "● Stalled jobs (worker died mid-run): automatically detected and re-queued",
    "● DLQ (Dead Letter Queue): jobs that failed MAX_RETRIES times — investigate root cause",
    "● Job statuses: PENDING / QUEUED / RUNNING / WAITING_APPROVAL / APPROVED / FAILED",
    "● AgentJob table: full payload, result, error details, and attempt count per job",
]
add_textbox_lines(sl,
                  Inches(0.4), Inches(1.45),
                  Inches(12.5), Inches(5.7),
                  lines, 14.5, color=WHITE, line_spacing=26)
add_slide_number(sl, 20)

# ===========================================================================
# SLIDE 21 — Log Monitoring & Alerts
# ===========================================================================
sl = prs.slides.add_slide(BLANK)
fill_slide_bg(sl)
add_header_bar(sl, "Log Monitoring & Alerts")

lines = [
    "● Railway Logs: real-time NestJS stdout — structured JSON logs, searchable",
    "● Search logs for: ERROR, WARN, specific userId, jobId, or agentName",
    "● Sentry: automatic error capture for both frontend (@sentry/nextjs) and backend",
    "● Prometheus metrics: GET /metrics endpoint, scraped by infra/monitoring stack",
    "● Grafana dashboards: http_request_duration_ms histogram, queue depth, error rate",
    "● Key alerts to configure: 500 error rate spike, queue depth > 100, DB failures",
    "● AgentLog table: every AI call logged — model, tokens in/out, cost, latency (ms)",
    "● Use AgentLog to track monthly AI spend and identify expensive agent operations",
    "● Bundle budget gate: 800KB per-route / 1500KB total JS — CI blocks on violation",
]
add_textbox_lines(sl,
                  Inches(0.4), Inches(1.45),
                  Inches(12.5), Inches(5.7),
                  lines, 14.5, color=WHITE, line_spacing=26)
add_slide_number(sl, 21)

# ===========================================================================
# SLIDE 22 — Database Operations
# ===========================================================================
sl = prs.slides.add_slide(BLANK)
fill_slide_bg(sl)
add_header_bar(sl, "Database Operations")

lines = [
    "● Access: Railway → PostgreSQL plugin → Connect → psql  (or Railway CLI)",
    "● Migrations: run  npx prisma migrate deploy  in API service shell before restart",
    "● NEVER alter the database directly — always go through Prisma migration files",
    "● Backups: Railway creates daily automated backups (Pro plan — highly recommended)",
    "● Manual backup: pg_dump via Railway CLI before any risky migration",
    "● Schema changes: update prisma/schema.prisma → generate migration → deploy",
    "● Monitoring: Railway Postgres metrics tab → check for slow queries over 1 second",
    "● Append-only tables: CreditLedger and AuditLog — rows are NEVER updated or deleted",
    "● Wallet balance: read-cache only; CreditLedger is the single source of truth",
]
add_textbox_lines(sl,
                  Inches(0.4), Inches(1.45),
                  Inches(12.5), Inches(5.7),
                  lines, 14.5, color=WHITE, line_spacing=26)
add_slide_number(sl, 22)

# ===========================================================================
# SLIDE 23 — Incident Response Checklist
# ===========================================================================
sl = prs.slides.add_slide(BLANK)
fill_slide_bg(sl)
add_header_bar(sl, "Incident Response Checklist")

lines = [
    "● Site down → check Vercel status page + Railway service health dashboard",
    "● API errors → Railway Logs → search 'ERROR' → check DB connection string",
    "● OTP not sending → verify RESEND_API_KEY (or SMTP_*) in Railway Variables",
    "● Jobs stuck → BullMQ Board → check Redis connection → restart worker service",
    "● DB connection error → Railway PostgreSQL → check max_connections limit",
    "● High AI costs → AgentLog table → filter by tokensOut → identify expensive calls",
    "● User locked out → /admin panel → reset OTP / grant temporary access",
    "● Chargeback/fraud → set Wallet.rechargesFrozen = true until review complete",
    "● Compliance bypass attempt → audit Semgrep CI results + AuditLog table",
]
add_textbox_lines(sl,
                  Inches(0.4), Inches(1.45),
                  Inches(12.5), Inches(5.7),
                  lines, 14.5, color=WHITE, line_spacing=26)
add_slide_number(sl, 23)

# ===========================================================================
# SLIDE 24 — Closing
# ===========================================================================
sl = prs.slides.add_slide(BLANK)
fill_slide_bg(sl)

# Gold bar at top
add_rect(sl, 0, 0, W, Emu(109728), GOLD)  # ~12px

# Main closing line
add_textbox(sl,
            Inches(0.5), Inches(2.5),
            Inches(12.33), Inches(0.9),
            "You run the platform.",
            44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Admin URL
add_textbox(sl,
            Inches(0.5), Inches(3.5),
            Inches(12.33), Inches(0.6),
            "admin panel: sozialzync.vercel.app/admin",
            22, bold=False, color=GOLD, align=PP_ALIGN.CENTER)

# Infra links
add_textbox(sl,
            Inches(0.5), Inches(4.3),
            Inches(12.33), Inches(0.5),
            "Railway dashboard: railway.app  |  Vercel: vercel.com",
            15, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

# GitHub
add_textbox(sl,
            Inches(0.5), Inches(5.0),
            Inches(12.33), Inches(0.45),
            "github.com/MoshayM/AI-Creatorforce",
            14, bold=False, color=SLATE, align=PP_ALIGN.CENTER)

add_slide_number(sl, 24)

# ===========================================================================
# Save
# ===========================================================================
out_path = r"D:\project\creatorforce-ai\CreatorForce-SuperAdmin-Guide.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slide count: {len(prs.slides)}")
