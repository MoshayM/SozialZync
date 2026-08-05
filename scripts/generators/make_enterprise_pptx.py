"""
AI CreatorForce — Enterprise Owner's Guide PPTX generator
Generates CreatorForce-Enterprise-Guide.pptx (22 slides)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── palette ─────────────────────────────────────────────────────────────────
AMBER      = RGBColor(0xD9, 0x77, 0x06)   # #D97706
GOLD       = RGBColor(0xFC, 0xD3, 0x4D)   # #FCD34D
DARK_TEXT  = RGBColor(0x78, 0x35, 0x0F)   # #78350F
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)   # #FFFFFF
BG         = RGBColor(0xFF, 0xFB, 0xEB)   # #FFFBEB
DARK_HDR   = RGBColor(0x92, 0x40, 0x0E)   # #92400E
CREAM      = RGBColor(0xFE, 0xF3, 0xC7)   # #FEF3C7
AMBER_LITE = RGBColor(0xFD, 0xE6, 0x8A)   # #FDE68A  alt row tint

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)

# ── helpers ──────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_rect(slide, x, y, w, h, fill_color, line=False):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if not line:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, font_size, bold=False,
             color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = False
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_wrapped_text(slide, text, x, y, w, h, font_size, bold=False,
                     color=DARK_TEXT, align=PP_ALIGN.LEFT, line_spacing=None):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def add_bullet_textbox(slide, lines, x, y, w, h, font_size=15,
                       color=DARK_TEXT, bullet_char="●"):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"{bullet_char} {line}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        p.space_after = Pt(4)
    return txb


def set_cell_bg(cell, hex_color):
    tc    = cell._tc
    tcPr  = tc.get_or_add_tcPr()
    # remove existing fills
    for old in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(old)
    solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
    srgbClr   = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', hex_color)


def set_cell_text(cell, text, bold=False, font_size=12,
                  color=RGBColor(0,0,0), align=PP_ALIGN.LEFT):
    tf = cell.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    # clear existing runs
    for run in p.runs:
        run.text = ''
    if p.runs:
        run = p.runs[0]
    else:
        run = p.add_run()
    run.text = text
    run.font.bold  = bold
    run.font.size  = Pt(font_size)
    run.font.color.rgb = color


def add_header_bar(slide, title):
    """Amber header bar + white title text."""
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), AMBER)
    add_text(slide, title,
             Inches(0.3), Inches(0.2), Inches(12.7), Inches(0.9),
             font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)


def add_slide_number(slide, num, total=22):
    add_text(slide, f"{num} / {total}",
             Inches(12.2), Inches(7.1), Inches(1.0), Inches(0.3),
             font_size=10, color=AMBER, align=PP_ALIGN.RIGHT)


def section_break(prs, text, num):
    slide = blank_slide(prs)
    # full amber background
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, AMBER)
    # gold accent line
    add_rect(slide, Inches(1.5), Inches(4.4), Inches(10.33), Inches(0.06), GOLD)
    # big white title
    add_text(slide, text,
             Inches(0.5), Inches(2.5), Inches(12.33), Inches(1.5),
             font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_slide_number(slide, num)
    return slide


def header_slide_with_bullets(prs, title, bullets, num):
    slide = blank_slide(prs)
    # light bg
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, BG)
    add_header_bar(slide, title)
    add_bullet_textbox(slide, bullets,
                       Inches(0.5), Inches(1.35), Inches(12.33), Inches(5.8),
                       font_size=15, color=DARK_TEXT)
    add_slide_number(slide, num)
    return slide


# ── SLIDE 1 — Title ──────────────────────────────────────────────────────────

def slide_01(prs):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, AMBER)

    add_text(slide, "AI CreatorForce",
             Inches(0.5), Inches(2.0), Inches(12.33), Inches(0.75),
             font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, "Enterprise Owner's Guide",
             Inches(0.5), Inches(3.0), Inches(12.33), Inches(0.6),
             font_size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, "Scale your content operation — one platform, your whole team",
             Inches(0.5), Inches(3.8), Inches(12.33), Inches(0.5),
             font_size=18, color=CREAM, align=PP_ALIGN.CENTER)

    add_text(slide, "★ ENTERPRISE PLAN ★",
             Inches(0.5), Inches(4.5), Inches(12.33), Inches(0.4),
             font_size=15, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    add_text(slide, "sozialzync.vercel.app",
             Inches(0.5), Inches(5.2), Inches(12.33), Inches(0.4),
             font_size=14, color=WHITE, align=PP_ALIGN.CENTER)

    add_slide_number(slide, 1)
    return slide


# ── SLIDE 4 — comparison table ───────────────────────────────────────────────

def slide_04(prs):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, BG)
    add_header_bar(slide, "Enterprise vs Pro vs Free — What's Different")

    cols = 4
    rows = 6
    tbl  = slide.shapes.add_table(
        rows, cols,
        Inches(0.4), Inches(1.35), Inches(12.53), Inches(5.7)
    ).table

    headers = ["Feature", "Free", "Pro", "Enterprise"]
    data = [
        ("Team seats",      "1",    "1",              "Unlimited"),
        ("Shared wallet",   "No",   "No",             "Yes — org-level"),
        ("Multi-channel",   "3",    "Unlimited",      "Unlimited + org view"),
        ("Analytics",       "Basic","Standard",       "Advanced + team view"),
        ("Support",         "Community","Email",      "Dedicated manager"),
    ]

    # header row
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        set_cell_bg(cell, "D97706")
        set_cell_text(cell, hdr, bold=True, font_size=13,
                      color=WHITE, align=PP_ALIGN.CENTER)

    # data rows
    row_colors = ["FFFBEB", "FDE68A"]
    for ri, row_data in enumerate(data):
        bg_hex = row_colors[ri % 2]
        for ci, val in enumerate(row_data):
            cell = tbl.cell(ri + 1, ci)
            set_cell_bg(cell, bg_hex)
            set_cell_text(cell, val, bold=(ci == 0), font_size=13,
                          color=DARK_TEXT, align=PP_ALIGN.CENTER)

    add_slide_number(slide, 4)
    return slide


# ── SLIDE 7 — roles table ─────────────────────────────────────────────────────

def slide_07(prs):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, BG)
    add_header_bar(slide, "Team Roles & Permissions")

    cols = 4
    rows = 6
    tbl  = slide.shapes.add_table(
        rows, cols,
        Inches(0.4), Inches(1.35), Inches(12.53), Inches(5.7)
    ).table

    headers = ["Role", "Can Create Content", "Can Publish", "Billing Access"]
    data = [
        ("MEMBER",        "Yes (own projects)",     "Needs approval",    "No"),
        ("TEAM_MANAGER",  "Yes + review others",    "Approve team content", "No"),
        ("BILLING_ADMIN", "No content",             "No",                "Yes — manage wallet"),
        ("ORG_ADMIN",     "Yes + full team",        "Full publish",      "Yes"),
        ("OWNER",         "Yes + org settings",     "Full publish",      "Full + plan change"),
    ]

    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        set_cell_bg(cell, "D97706")
        set_cell_text(cell, hdr, bold=True, font_size=13,
                      color=WHITE, align=PP_ALIGN.CENTER)

    row_colors = ["FFFBEB", "FDE68A"]
    for ri, row_data in enumerate(data):
        bg_hex = row_colors[ri % 2]
        for ci, val in enumerate(row_data):
            cell = tbl.cell(ri + 1, ci)
            set_cell_bg(cell, bg_hex)
            set_cell_text(cell, val, bold=(ci == 0), font_size=13,
                          color=DARK_TEXT, align=PP_ALIGN.CENTER)

    add_slide_number(slide, 7)
    return slide


# ── SLIDE 20 — two-column ROI ─────────────────────────────────────────────────

def slide_20(prs):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, BG)
    add_header_bar(slide, "Cost Savings vs. Traditional Team")

    # divider line
    add_rect(slide, Inches(6.7), Inches(1.4), Inches(0.03), Inches(5.8), AMBER)

    left_lines = [
        "Traditional content team (monthly):",
        "",
        "● Research analyst: $3,000–5,000/mo",
        "● Scriptwriter: $3,000–6,000/mo",
        "● SEO specialist: $2,000–4,000/mo",
        "● Compliance reviewer: $2,500–4,000/mo",
        "",
        "● Total: $10,500–19,000/mo per channel",
    ]
    right_lines = [
        "AI CreatorForce Enterprise:",
        "",
        "● Enterprise plan: $199/mo (platform)",
        "● AI credits for 20 videos: ~$50–200/mo",
        "● Total: ~$250–400/mo per channel",
        "",
        "● Savings: 95%+ cost reduction",
        "● Speed: 2-4 hours per video vs 2-3 days",
    ]

    def col_box(slide, lines, x):
        txb = slide.shapes.add_textbox(x, Inches(1.4), Inches(5.9), Inches(5.8))
        tf  = txb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = line
            if line.startswith("Traditional") or line.startswith("AI Creator"):
                run.font.bold = True
                run.font.size = Pt(16)
                run.font.color.rgb = DARK_HDR
            elif line.startswith("● Total") or line.startswith("● Savings") or line.startswith("● Speed"):
                run.font.bold  = True
                run.font.size  = Pt(14)
                run.font.color.rgb = DARK_TEXT
            elif line == "":
                run.font.size = Pt(8)
                run.font.color.rgb = DARK_TEXT
            else:
                run.font.size = Pt(14)
                run.font.color.rgb = DARK_TEXT
            p.space_after = Pt(2)

    col_box(slide, left_lines,  Inches(0.4))
    col_box(slide, right_lines, Inches(6.9))

    add_slide_number(slide, 20)
    return slide


# ── SLIDE 22 — Closing ────────────────────────────────────────────────────────

def slide_22(prs):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, AMBER)

    add_text(slide,
             "Scale your content. Grow your channels. Build your empire.",
             Inches(0.5), Inches(2.3), Inches(12.33), Inches(0.8),
             font_size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, "Get started at: sozialzync.vercel.app",
             Inches(0.5), Inches(3.3), Inches(12.33), Inches(0.6),
             font_size=24, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, "Questions? Contact your dedicated account manager",
             Inches(0.5), Inches(4.1), Inches(12.33), Inches(0.5),
             font_size=18, color=CREAM, align=PP_ALIGN.CENTER)

    add_text(slide, "AI CreatorForce — Enterprise Content Operations",
             Inches(0.5), Inches(5.0), Inches(12.33), Inches(0.4),
             font_size=14, color=WHITE, align=PP_ALIGN.CENTER)

    add_slide_number(slide, 22)
    return slide


# ── BUILD ALL 22 SLIDES ───────────────────────────────────────────────────────

def build():
    prs = new_prs()

    # Slide 1 — Title
    slide_01(prs)

    # Slide 2 — Section break
    section_break(prs, "WHAT YOU GET WITH ENTERPRISE", 2)

    # Slide 3 — Enterprise at a Glance
    header_slide_with_bullets(prs, "Enterprise Plan — At a Glance", [
        "Unlimited projects and AI agent usage across your entire team",
        "All 15 AI agents with priority processing queue",
        "Shared organization wallet — one billing account for all team members",
        "Team member seats with role-based permissions",
        "Multiple YouTube channels under one organization",
        "White-label dashboard: your brand, your team's workspace",
        "Dedicated account manager + priority support SLA",
        "Full compliance and audit trail for every piece of content",
    ], 3)

    # Slide 4 — Comparison table
    slide_04(prs)

    # Slide 5 — Section break
    section_break(prs, "ORGANIZATION SETUP", 5)

    # Slide 6 — Setting Up Your Organization
    header_slide_with_bullets(prs, "Setting Up Your Organization", [
        "Navigate to: App → Team Workspaces (requires Enterprise plan)",
        "Click \"Create Organization\" — enter org name and billing email",
        "Set up shared org wallet: load credits for the team to share",
        "Add YouTube channels: each channel gets its own brand profile",
        "Invite team members via email — they receive OTP login link",
        "Assign roles before they start working",
    ], 6)

    # Slide 7 — Roles table
    slide_07(prs)

    # Slide 8 — Managing Team Members
    header_slide_with_bullets(prs, "Managing Team Members", [
        "Invite: Organization Settings → Members → Invite by Email",
        "Set role at invite time — can be changed later by ORG_ADMIN+",
        "Remove member: immediately revokes access, content stays in org",
        "View member activity: projects created, AI usage, credits consumed",
        "Budget per member: set individual credit caps to prevent overspend",
        "Members see only their assigned channels and projects (RBAC enforced)",
    ], 8)

    # Slide 9 — Channel Management at Scale
    header_slide_with_bullets(prs, "Channel Management at Scale", [
        "Add multiple YouTube channels to your organization",
        "Each channel has its own: niche profile, brand voice, tone settings",
        "Assign channels to specific team members or teams",
        "Cross-channel analytics: compare performance across all your channels",
        "Channel-level budget caps: limit AI spend per channel per month",
        "Bulk operations: apply SEO templates across multiple channels",
        "Brand kit: logo, color palette, intro/outro styles shared per channel",
    ], 9)

    # Slide 10 — Section break
    section_break(prs, "BILLING & COST CONTROL", 10)

    # Slide 11 — Org Wallet
    header_slide_with_bullets(prs, "Organization Wallet — How It Works", [
        "Single shared wallet holds credits for the entire organization",
        "Credits are deducted when AI agents run (per-token pricing)",
        "Credit types consumed in order: Trial → Purchased → Bonus → Referral",
        "Oldest expiry lot consumed first (FIFO by expiry date)",
        "Wallet owner: BILLING_ADMIN or ORG_ADMIN tops up credits",
        "Real-time balance visible to all admins in: App → Wallet",
    ], 11)

    # Slide 12 — Budget Controls
    header_slide_with_bullets(prs, "Budget Controls & Caps", [
        "Organization monthly spend cap: hard limit — agents stop when reached",
        "Per-channel monthly cap: limit how much one channel consumes",
        "Per-member cap: limit individual creator's daily/monthly AI usage",
        "Low balance alert: notification at configurable threshold (e.g. < 5,000 credits)",
        "Budget report: weekly email to BILLING_ADMIN with spend breakdown",
        "Overspend protection: jobs queued but not started when cap would be exceeded",
    ], 12)

    # Slide 13 — Credit Costs
    header_slide_with_bullets(prs, "Understanding Credit Costs", [
        "Credits track AI token usage: ~1 credit per 1,000 tokens",
        "Script generation (full video): ~800–1,200 credits",
        "Fact check pass: ~200–400 credits",
        "Compliance check: ~150–300 credits",
        "SEO optimization: ~200–400 credits",
        "Copilot query: ~50–150 credits",
        "Shorts clip analysis (per video): ~300–600 credits",
        "Note: actual costs vary by model tier and content length",
    ], 13)

    # Slide 14 — Section break
    section_break(prs, "CONTENT OPERATIONS AT SCALE", 14)

    # Slide 15 — Multi-Channel Pipeline
    header_slide_with_bullets(prs, "Multi-Channel Content Pipeline", [
        "Each channel runs its own independent content pipeline",
        "Pipelines can run in parallel — no bottleneck between channels",
        "Channel A: research in progress while Channel B is in compliance review",
        "SupervisorAgent manages priority queue: Enterprise = priority processing",
        "Team manager can view all active pipelines on the team dashboard",
        "Assign specific AI agents to specific channels (e.g. premium model for main channel)",
    ], 15)

    # Slide 16 — Content Calendar
    header_slide_with_bullets(prs, "Content Calendar & Publishing Schedule", [
        "Set publishing cadence per channel: daily / 3x week / weekly / custom",
        "Content Calendar view: drag-and-drop scheduled videos per channel",
        "Auto-suggest optimal publish times based on audience analytics",
        "Scheduled publish: enabled per-channel by ORG_ADMIN",
        "Even with auto-schedule enabled: EVERY video must pass compliance first",
        "Human approval still required unless auto-approve is explicitly enabled",
        "Bulk approve: team manager reviews + approves multiple videos at once",
    ], 16)

    # Slide 17 — Compliance at Scale
    header_slide_with_bullets(prs, "Compliance at Enterprise Scale", [
        "ComplianceAgent runs on EVERY piece of content — no exceptions",
        "Enterprise benefit: extended compliance report with risk category breakdown",
        "Monetization Safety Score: must exceed org-configured threshold",
        "Compliance dashboard: org-wide pass/fail rate, common failure reasons",
        "Fact-check citations: every claim has a source — full audit trail",
        "Content blocked by compliance: notification to creator + manager",
        "Audit log: complete history of compliance decisions, downloadable CSV",
    ], 17)

    # Slide 18 — Analytics
    header_slide_with_bullets(prs, "Analytics — Team & Channel Performance", [
        "Individual channel analytics: views, watch time, CTR, subscribers",
        "Team-level analytics: aggregated across all org channels",
        "Content performance by agent: which AI-generated scripts get best retention?",
        "Top performing creators on your team: ranked by channel growth",
        "AI cost efficiency: cost per published video, cost per 1K views",
        "Weekly performance report: PDF summary to org admins via email",
        "Custom date ranges, channel filters, content-type filters",
    ], 18)

    # Slide 19 — Section break
    section_break(prs, "ENTERPRISE ROI", 19)

    # Slide 20 — Cost Savings two-column
    slide_20(prs)

    # Slide 21 — Getting Maximum Value
    header_slide_with_bullets(prs, "Getting Maximum Value from Enterprise", [
        "Assign dedicated channels to dedicated team members — clear ownership",
        "Set up brand kits per channel — AI adapts style automatically",
        "Use Shorts Studio on your top 5 existing videos immediately — free reach",
        "Enable the content calendar — consistent posting = algorithm boost",
        "Review the analytics weekly as a team — data-driven topic selection",
        "Set conservative budget caps first — increase once team is calibrated",
        "Use the compliance dashboard to identify and fix recurring content issues",
    ], 21)

    # Slide 22 — Closing
    slide_22(prs)

    out = r"D:\project\creatorforce-ai\CreatorForce-Enterprise-Guide.pptx"
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Slide count: {len(prs.slides)}")


if __name__ == "__main__":
    build()
