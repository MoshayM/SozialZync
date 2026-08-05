"""
make_dev_pptx.py — AI CreatorForce Developer's Reference PowerPoint Generator
Creates CreatorForce-Developer-Guide.pptx (28 slides)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ──────────────────────────────────────────────────────────────────
TEAL        = RGBColor(0x08, 0x91, 0xB2)   # #0891B2
ACCENT      = RGBColor(0x22, 0xD3, 0xEE)   # #22D3EE
DARK        = RGBColor(0x0C, 0x4A, 0x6E)   # #0C4A6E
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)   # #FFFFFF
LIGHT_BG    = RGBColor(0xF0, 0xF9, 0xFF)   # #F0F9FF
CODE_BG     = RGBColor(0x1E, 0x29, 0x3B)   # #1E293B
GRAY        = RGBColor(0x94, 0xA3, 0xB8)   # slate-400
BAE6FD      = RGBColor(0xBA, 0xE6, 0xFD)   # #BAE6FD

SW = Inches(13.33)
SH = Inches(7.5)

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH
blank_layout = prs.slide_layouts[6]  # completely blank

# ── Helper functions ──────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_color, border=False):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if not border:
        shape.line.fill.background()
    return shape


def add_textbox(slide, x, y, w, h, text, font_size=15, bold=False,
                color=DARK, align=PP_ALIGN.LEFT, font_name="Calibri",
                word_wrap=True, line_spacing=None):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    if line_spacing:
        from pptx.util import Pt as PT
        from pptx.oxml.ns import qn
        from lxml import etree
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
        spcPts.set('val', str(int(line_spacing * 100)))
    return txBox


def add_multiline_textbox(slide, x, y, w, h, lines, font_size=15, bold=False,
                          color=DARK, align=PP_ALIGN.LEFT, font_name="Calibri",
                          line_spacing_pt=None):
    """Add a textbox with multiple lines, each as a paragraph."""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
        if line_spacing_pt:
            from pptx.oxml.ns import qn
            from lxml import etree
            pPr = p._p.get_or_add_pPr()
            lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
            spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
            spcPts.set('val', str(int(line_spacing_pt * 100)))
    return txBox


def add_header(slide, title):
    """Add teal header bar + white title text."""
    add_rect(slide, Inches(0), Inches(0), SW, Inches(1.2), TEAL)
    add_textbox(slide, Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.9),
                title, font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)


def add_slide_number(slide, num):
    add_textbox(slide, Inches(12.5), Inches(7.1), Inches(0.8), Inches(0.3),
                str(num), font_size=11, color=GRAY, align=PP_ALIGN.RIGHT)


def add_code_block(slide, x, y, w, h, lines, font_size=12):
    """Dark rectangle with white monospace text."""
    add_rect(slide, x, y, w, h, CODE_BG)
    code_text = "\n".join(lines)
    txBox = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1),
                                      w - Inches(0.3), h - Inches(0.2))
    tf = txBox.text_frame
    tf.word_wrap = False
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = WHITE
        run.font.name = "Courier New"


def body_start():
    return Inches(1.4)


# ── Slide 1: Title ────────────────────────────────────────────────────────────
def slide1(prs):
    s = prs.slides.add_slide(blank_layout)
    add_rect(s, 0, 0, SW, SH, TEAL)
    add_textbox(s, Inches(1), Inches(2.3), Inches(11.33), Inches(0.9),
                "AI CreatorForce", font_size=52, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(3.3), Inches(11.33), Inches(0.6),
                "Technical Deep Dive — Developer's Reference", font_size=22,
                bold=False, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(4.2), Inches(11.33), Inches(0.5),
                "NestJS · Next.js · TypeScript · PostgreSQL · BullMQ · 15 AI Agents",
                font_size=15, color=BAE6FD, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(5.2), Inches(11.33), Inches(0.4),
                "sozialzync.vercel.app | github.com/MoshayM/AI-Creatorforce",
                font_size=13, color=WHITE, align=PP_ALIGN.CENTER)
    add_slide_number(s, 1)


# ── Slide 2: Section break ────────────────────────────────────────────────────
def section_break(prs, title, num):
    s = prs.slides.add_slide(blank_layout)
    add_rect(s, 0, 0, SW, SH, TEAL)
    add_textbox(s, Inches(1), Inches(2.9), Inches(11.33), Inches(1.2),
                title, font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_slide_number(s, num)


# ── Slide 3: What AI CreatorForce Does ────────────────────────────────────────
def slide3(prs):
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "What AI CreatorForce Does")
    lines = [
        "● SaaS platform: AI-powered YouTube Content Operating System",
        "● Multi-tenant: users, organizations, channels, projects",
        "● 15 specialized AI agents orchestrated by SupervisorWorker",
        "● Full content pipeline: research → script → compliance → publish",
        "● Real-time job progress via WebSockets + BullMQ queues",
        "● Auth: JWT + OTP (email/SMS) + Google OAuth",
    ]
    add_multiline_textbox(s, Inches(0.4), body_start(), Inches(12.5), Inches(5.8),
                          lines, font_size=17, color=DARK, line_spacing_pt=26)
    add_slide_number(s, 3)


# ── Slide 4: Repository Structure ─────────────────────────────────────────────
def slide4(prs):
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Repository Structure (Monorepo)")
    tree = [
        "creatorforce-ai/",
        "├── apps/",
        "│   ├── web/          ← Next.js 15 frontend (Vercel)",
        "│   └── api/          ← NestJS backend (Railway)",
        "├── packages/",
        "│   ├── agents/       ← AI agent implementations",
        "│   ├── shared/       ← Zod schemas, types, utils",
        "│   ├── prompts/      ← Versioned prompt templates",
        "│   └── config/       ← ESLint, TSConfig, Tailwind presets",
        "├── infra/            ← Docker, IaC, GitHub Actions",
        "├── n8n/              ← Exported workflow definitions",
        "└── CLAUDE.md         ← AI agent operating contract",
    ]
    add_code_block(s, Inches(0.4), body_start(), Inches(12.5), Inches(5.7),
                   tree, font_size=14)
    add_slide_number(s, 4)


# ── Slide 5: Tech Stack Frontend ──────────────────────────────────────────────
def slide5(prs):
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Tech Stack: Frontend (apps/web)")
    lines = [
        "● Framework: Next.js 15 with App Router",
        "● Language: TypeScript strict mode (no any without // @reason:)",
        "● Styling: Tailwind CSS + custom design tokens",
        "● State/fetching: TanStack Query v5 (React Query)",
        "● UI primitives: Radix UI / Lucide icons",
        "● Mobile: bottom nav, sidebar drawer, viewport/iOS fixes",
        "● Deploy: Vercel (auto-deploy on git push to master)",
    ]
    add_multiline_textbox(s, Inches(0.4), body_start(), Inches(12.5), Inches(5.8),
                          lines, font_size=17, color=DARK, line_spacing_pt=26)
    add_slide_number(s, 5)


# ── Slide 6: Tech Stack Backend ───────────────────────────────────────────────
def slide6(prs):
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Tech Stack: Backend (apps/api)")
    lines = [
        "● Framework: NestJS 10 with modular architecture",
        "● ORM: Prisma 6 (type-safe, schema-first, migration runner)",
        "● Database: PostgreSQL 16 (Railway plugin)",
        "● Queue: BullMQ 5 + Redis 7 (Railway plugin) — async jobs",
        "● WebSockets: NestJS Gateway (socket.io) for real-time job events",
        "● Auth: JWT access + refresh tokens, Passport strategies",
        "● Deploy: Railway (persistent workers, persistent connections)",
    ]
    add_multiline_textbox(s, Inches(0.4), body_start(), Inches(12.5), Inches(5.8),
                          lines, font_size=17, color=DARK, line_spacing_pt=26)
    add_slide_number(s, 6)


# ── Slide 7: Tech Stack AI & External ────────────────────────────────────────
def slide7(prs):
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Tech Stack: AI & External Services")
    lines = [
        "● AI providers: Anthropic Claude (primary), OpenAI (fallback), Gemini",
        "● Video generation: Veo, Kling, Runway, Pika, Luma",
        "● Music generation: Suno, Udio, Stable Audio",
        "● Email OTP: Resend (primary) → SMTP (fallback)",
        "● SMS OTP: Twilio",
        "● Payments: Stripe 17",
        "● Monitoring: Sentry (errors), prom-client/Prometheus (metrics), BullMQ Dashboard",
    ]
    add_multiline_textbox(s, Inches(0.4), body_start(), Inches(12.5), Inches(5.8),
                          lines, font_size=17, color=DARK, line_spacing_pt=26)
    add_slide_number(s, 7)


# ── Slide 9: NestJS Module Structure ──────────────────────────────────────────
def slide9(prs):
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "NestJS Module Structure")
    lines = [
        "● One NestJS module per Core Engine — mirrors the product domain",
        "● Modules: AuthModule, ChannelsModule, ProjectsModule, JobsModule",
        "● Modules: WalletModule, BillingModule, OrgsModule, PublishingModule",
        "● Modules: AnalyticsModule, AgentsModule, NotificationsModule, HealthModule",
        "● Each module: Controller → Service → Repository → Prisma",
        "● Shared services: PrismaService, RedisService, AiClientService",
        "● Guards applied globally: JwtAuthGuard, MetricsInterceptor",
    ]
    add_multiline_textbox(s, Inches(0.4), body_start(), Inches(12.5), Inches(5.8),
                          lines, font_size=17, color=DARK, line_spacing_pt=26)
    add_slide_number(s, 9)


# ── Slide 10: BullMQ Async Job Architecture ───────────────────────────────────
def slide10(prs):
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "BullMQ Async Job Architecture")
    lines = [
        "● Rule: anything > 2 seconds runs as a BullMQ job, NEVER inline",
        "● Single logical queue: AGENT_QUEUE (all agent + media jobs)",
        "● Job types: ~50 JobType enum values covering every pipeline step",
        "● Job states: PENDING → QUEUED → RUNNING → COMPLETED / FAILED",
        "● Failed jobs: retry with exponential backoff → QualityControlAgent",
        "● Real-time progress: Socket.io events on each job state change",
        "● Dashboard: BullMQ Board UI for queue monitoring",
    ]
    add_multiline_textbox(s, Inches(0.4), body_start(), Inches(12.5), Inches(5.8),
                          lines, font_size=17, color=DARK, line_spacing_pt=26)
    add_slide_number(s, 10)


# ── Slide 11: Database Schema ─────────────────────────────────────────────────
def slide11(prs):
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Database Schema (Key Tables)")

    left_lines = [
        "Core entities:",
        "● User (id, email, phone, role, plan)",
        "● Subscription (userId, tier, stripeId)",
        "● Channel (userId, ytChannelId, brandProfile)",
        "● Project (channelId, contentType, status)",
        "● AgentJob (projectId, type, status, result)",
        "● AgentLog (jobId, model, tokens, latency)",
    ]
    right_lines = [
        "Financial & compliance:",
        "● Wallet (userId, balanceCredits)",
        "● CreditLedger (walletId, type, amount)",
        "● ComplianceResult (jobId, passed, score)",
        "● Approval (projectId, jobId, status)",
        "● AuditLog (userId, action, meta) — append-only",
        "● AuthSession (userId, refreshToken, rotated)",
    ]

    add_multiline_textbox(s, Inches(0.4), body_start(), Inches(6.2), Inches(5.8),
                          left_lines, font_size=15, color=DARK, line_spacing_pt=22)
    add_multiline_textbox(s, Inches(6.8), body_start(), Inches(6.1), Inches(5.8),
                          right_lines, font_size=15, color=DARK, line_spacing_pt=22)
    # vertical divider
    add_rect(s, Inches(6.6), body_start(), Inches(0.03), Inches(5.6), ACCENT)
    add_slide_number(s, 11)


# ── Slide 12: Authentication & Security ───────────────────────────────────────
def slide12(prs):
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Authentication & Security")
    lines = [
        "● JWT: short-lived access token (15m) + long-lived refresh token (7d, rotating)",
        "● OTP login: 6-digit code → Resend/SMTP/Twilio → verify endpoint",
        "● Google OAuth via Passport GoogleStrategy; tokens encrypted at rest (AES-256-GCM)",
        "● Role hierarchy: MEMBER < OWNER < SUPER_ADMIN",
        "● Guards: JwtAuthGuard, OwnerGuard, PermissionsGuard on every protected route",
        "● Secrets: all via environment variables — zero secrets in code",
        "● Production startup guard: refuses to boot without JWT_SECRET + TOKEN_ENCRYPTION_KEY",
    ]
    add_multiline_textbox(s, Inches(0.4), body_start(), Inches(12.5), Inches(5.8),
                          lines, font_size=16, color=DARK, line_spacing_pt=25)
    add_slide_number(s, 12)


# ── Slide 14: Next.js App Router Structure ────────────────────────────────────
def slide14(prs):
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Next.js App Router Structure")
    lines = [
        "● app/(auth)/       — login, register, OTP verify pages",
        "● app/(dash)/       — main app: projects, channels, wallet, analytics, orgs",
        "● app/admin/        — Super Admin panel (role-gated)",
        "● app/api/          — Next.js API routes (proxies to NestJS backend)",
        "● components/       — shared UI: PlanGate, DevicePreview, BottomNav",
        "● lib/              — api.ts (typed API client), plan.ts (usePlan hook), auth.ts",
    ]
    add_multiline_textbox(s, Inches(0.4), body_start(), Inches(12.5), Inches(5.5),
                          lines, font_size=17, color=DARK, line_spacing_pt=26)
    add_slide_number(s, 14)


# ── Slide 15: SSR vs Client Components ───────────────────────────────────────
def slide15(prs):
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "SSR vs Client Components")
    lines = [
        "● Server Components by default — no JavaScript sent to client",
        "● Client Components only when interactive (marked 'use client')",
        "● IMPORTANT: 'use client' components still SSR on first render!",
        "● Rule: NEVER read localStorage at render time — causes hydration mismatch",
        "● Pattern: useState(defaultValue) + useEffect(() => read localStorage)",
        "● Example fix: contentType state in /projects/[id]/page.tsx (commit b5bb92c)",
    ]
    add_multiline_textbox(s, Inches(0.4), body_start(), Inches(12.5), Inches(5.5),
                          lines, font_size=17, color=DARK, line_spacing_pt=26)
    add_slide_number(s, 15)


# ── Slide 16: Plan Gating System ─────────────────────────────────────────────
def slide16(prs):
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Plan Gating System (Frontend)")
    code_lines = [
        "// plan-gate.tsx — UI affordance only; API always enforces server-side",
        "",
        "Plan tiers: FREE(0) < STARTER(1) < PRO(2) < ENTERPRISE(3) < AGENCY(4)",
        "",
        "usePlanGate()  → reads JWT plan field after mount (useEffect, avoids SSR mismatch)",
        "useIsAdmin()   → reads JWT role field; returns true for SUPER_ADMIN / OWNER",
        "",
        "PlanGate component logic:",
        "  const allowed = isAdmin || planAtLeast(userPlan, requiredPlan)",
        "  if (allowed) return children",
        "  else         return <LockOverlay upgradeLink='/wallet' />",
    ]
    add_code_block(s, Inches(0.4), body_start(), Inches(12.5), Inches(5.7),
                   code_lines, font_size=13)
    add_slide_number(s, 16)


# ── Slide 18: Agent Architecture Principles ───────────────────────────────────
def slide18(prs):
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Agent Architecture Principles")
    lines = [
        "● SupervisorWorker orchestrates — individual agents are stateless + idempotent",
        "● Agent inputs/outputs validated against Zod schemas (reject + retry on failure)",
        "● Prompts stored in packages/prompts (versioned, never inlined in code)",
        "● All provider calls via shared aiClient: retry, fallback provider, token accounting",
        "● On schema failure: retry up to MAX_AGENT_RETRIES → route to QualityControlAgent",
        "● Every agent call emits: {agentName, model, tokens, latencyMs, cost} trace event",
        "● Stored in AgentLog table for cost monitoring and debugging",
    ]
    add_multiline_textbox(s, Inches(0.4), body_start(), Inches(12.5), Inches(5.8),
                          lines, font_size=17, color=DARK, line_spacing_pt=26)
    add_slide_number(s, 18)


# ── Slide 19: The 15 AI Agents ────────────────────────────────────────────────
def slide19(prs):
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "The 15+ AI Agents")

    left_lines = [
        "● SupervisorWorker — orchestrates pipeline",
        "● ResearchAgent — trends, keywords, competitors",
        "● AudienceAgent — viewer profiling",
        "● ScriptAgent — full script generation",
        "● FactCheckAgent — claim verification + citations",
        "● ComplianceAgent — monetization safety (HARD GATE)",
        "● SEOAgent — title, description, tags",
        "● TrendAgent — trend discovery + signals",
    ]
    right_lines = [
        "● MetadataAgent — cards, chapters, pinned comment",
        "● ImageAgent — thumbnail/asset brief",
        "● StoryboardAgent — scene-by-scene visual plan",
        "● VoiceAgent — TTS script + pacing notes",
        "● MusicAgent — background score brief",
        "● VideoAgent — B-roll/generation prompts",
        "● EditPlanAgent — edit commands for editor",
        "● QualityControlAgent — catch + retry failures",
    ]

    add_multiline_textbox(s, Inches(0.4), body_start(), Inches(6.2), Inches(5.8),
                          left_lines, font_size=15, color=DARK, line_spacing_pt=23)
    add_multiline_textbox(s, Inches(6.8), body_start(), Inches(6.1), Inches(5.8),
                          right_lines, font_size=15, color=DARK, line_spacing_pt=23)
    add_rect(s, Inches(6.6), body_start(), Inches(0.03), Inches(5.6), ACCENT)
    add_slide_number(s, 19)


# ── Slide 20: Adding a New Agent ──────────────────────────────────────────────
def slide20(prs):
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Adding a New Agent — Checklist")
    lines = [
        "● 1. Define InputSchema + OutputSchema in packages/shared (Zod)",
        "● 2. Create prompt template in packages/prompts/{agentName}/v1.ts",
        "● 3. Extend BaseAgent<TInput, TOutput> — implement name, systemPrompt, run()",
        "● 4. Use callStructured() (not callAI()) for stored/downstream responses",
        "● 5. Validate response against OutputSchema; retry on failure up to MAX_AGENT_RETRIES",
        "● 6. Register new JobType in shared JobType enum + SupervisorWorker handler",
        "● 7. Emit trace event (handled by BaseAgent if using shared call methods)",
        "● 8. Write co-located unit test: {agentName}.agent.spec.ts",
    ]
    add_multiline_textbox(s, Inches(0.4), body_start(), Inches(12.5), Inches(5.8),
                          lines, font_size=16, color=DARK, line_spacing_pt=24)
    add_slide_number(s, 20)


# ── Slide 22: Deployment Architecture ────────────────────────────────────────
def slide22(prs):
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Deployment Architecture")

    left_lines = [
        "Vercel (Frontend):",
        "● Next.js app auto-deploys on git push",
        "● Preview URL for every PR branch",
        "● Global CDN — zero cold start for pages",
        "● Env var: NEXT_PUBLIC_API_URL → Railway",
        "● Bundle budget gate: 800 KB / route max",
    ]
    right_lines = [
        "Railway (Backend):",
        "● NestJS API — persistent process",
        "● BullMQ worker — persistent background worker",
        "● PostgreSQL plugin — managed DB",
        "● Redis plugin — managed cache/queue",
        "● WebSocket server — persistent connections",
    ]

    add_multiline_textbox(s, Inches(0.4), body_start(), Inches(6.2), Inches(5.8),
                          left_lines, font_size=16, color=DARK, line_spacing_pt=24)
    add_multiline_textbox(s, Inches(6.8), body_start(), Inches(6.1), Inches(5.8),
                          right_lines, font_size=16, color=DARK, line_spacing_pt=24)
    add_rect(s, Inches(6.6), body_start(), Inches(0.03), Inches(5.6), ACCENT)
    add_slide_number(s, 22)


# ── Slide 23: Environment Variables ──────────────────────────────────────────
def slide23(prs):
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Environment Variables")

    left_code = [
        "Vercel (frontend):",
        "NEXT_PUBLIC_API_URL=https://api.railway.app",
    ]
    right_code = [
        "Railway (backend — required):",
        "DATABASE_URL, REDIS_URL",
        "JWT_SECRET, JWT_REFRESH_SECRET",
        "TOKEN_ENCRYPTION_KEY (min 32 chars)",
        "ANTHROPIC_API_KEY, OPENAI_API_KEY, GEMINI_API_KEY",
        "RESEND_API_KEY, RESEND_FROM  (email OTP)",
        "TWILIO_ACCOUNT_SID, TWILIO_AUTH_TOKEN, TWILIO_FROM",
        "STRIPE_SECRET_KEY, STRIPE_WEBHOOK_SECRET",
        "SUPER_ADMIN_EMAILS, OWNER_EMAILS",
        "NODE_ENV=production",
        "CORS_ORIGIN=https://sozialzync.vercel.app",
    ]

    add_code_block(s, Inches(0.4), body_start(), Inches(5.8), Inches(1.5),
                   left_code, font_size=13)
    add_code_block(s, Inches(6.6), body_start(), Inches(6.5), Inches(5.5),
                   right_code, font_size=13)
    add_slide_number(s, 23)


# ── Slide 24: CI/CD Pipeline ──────────────────────────────────────────────────
def slide24(prs):
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "CI/CD Pipeline (GitHub Actions)")
    lines = [
        "● Trigger: push to master / pull_request targeting master",
        "● Step 1: lint — ESLint 10 across all packages and apps",
        "● Step 2: typecheck — tsc --noEmit (strict mode) + prisma generate",
        "● Step 3: unit-tests — Jest 29 with coverage collection",
        "● Step 4: build — turbo build (both apps/web and apps/api)",
        "● Step 5: security — pnpm audit --audit-level=high + dependency-review-action",
        "● Step 6: semgrep — SAST scan (ERROR severity blocks CI)",
        "● Step 7: zap-baseline — OWASP ZAP passive scan (HIGH findings block CI)",
        "● Step 8: e2e — Playwright cross-browser: chromium, firefox, webkit",
    ]
    add_multiline_textbox(s, Inches(0.4), body_start(), Inches(12.5), Inches(5.8),
                          lines, font_size=16, color=DARK, line_spacing_pt=23)
    add_slide_number(s, 24)


# ── Slide 25: Testing Strategy ────────────────────────────────────────────────
def slide25(prs):
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Testing Strategy")
    lines = [
        "● Unit tests: co-locate *.spec.ts next to each source file (apps/api/src)",
        "● Integration tests: required for all pipeline flows (agent → job → db)",
        "● Test philosophy: test contracts, schemas, and gates — not AI text output",
        "● No mocking the DB in integration tests (real DB via Postgres Docker service)",
        "● Agent tests: mock aiClient, test Zod validation + retry logic",
        "● E2E: Playwright for critical paths (login → create project → approve → publish)",
        "● Non-negotiable: ComplianceService, PublishingService, WalletService hard-cap",
    ]
    add_multiline_textbox(s, Inches(0.4), body_start(), Inches(12.5), Inches(5.8),
                          lines, font_size=16, color=DARK, line_spacing_pt=25)
    add_slide_number(s, 25)


# ── Slide 26: Golden Rules for Contributors ───────────────────────────────────
def slide26(prs):
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Golden Rules for Contributors")
    lines = [
        "● ComplianceAgent is a HARD GATE — never bypass it in any code path",
        "● Human-in-the-loop on publish — never auto-publish without explicit approval",
        "● No fabricated facts — every claim must trace to FactCheckAgent source",
        "● Secrets ONLY in env vars — never hardcode, never commit .env",
        "● TypeScript strict — no `any` without a // @reason: comment",
        "● Zod at every boundary — API input, agent output, env vars",
        "● Async work > 2s → BullMQ job, never inline in request handler",
    ]
    add_multiline_textbox(s, Inches(0.4), body_start(), Inches(12.5), Inches(5.8),
                          lines, font_size=17, color=DARK, line_spacing_pt=26)
    add_slide_number(s, 26)


# ── Slide 27: Common Pitfalls & Fixes ─────────────────────────────────────────
def slide27(prs):
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "Common Pitfalls & Fixes")
    lines = [
        "● SSR hydration mismatch → never read localStorage at render time (use useEffect)",
        "● CSS overflow clipping dropdowns → move absolute-positioned elements outside overflow:auto",
        "● Renamed export breaks imports → always add backward-compat alias",
        "● OTP silent failure in prod → check NODE_ENV before falling to dev console log",
        "● Missing Railway env var → real error now thrown instead of false success",
        "● BullMQ not processing → check Redis connection URL in Railway variables",
        "● prisma generate missing → run before typecheck, test, or build in CI",
    ]
    add_multiline_textbox(s, Inches(0.4), body_start(), Inches(12.5), Inches(5.8),
                          lines, font_size=16, color=DARK, line_spacing_pt=25)
    add_slide_number(s, 27)


# ── Slide 28: Closing ─────────────────────────────────────────────────────────
def slide28(prs):
    s = prs.slides.add_slide(blank_layout)
    add_rect(s, 0, 0, SW, SH, TEAL)
    add_textbox(s, Inches(1), Inches(2.8), Inches(11.33), Inches(0.9),
                "Build something remarkable.", font_size=40, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(3.8), Inches(11.33), Inches(0.6),
                "github.com/MoshayM/AI-Creatorforce", font_size=22,
                color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(4.6), Inches(11.33), Inches(0.5),
                "Docs: /docs directory in repo | App: sozialzync.vercel.app",
                font_size=15, color=BAE6FD, align=PP_ALIGN.CENTER)
    add_slide_number(s, 28)


# ── Build all 28 slides ───────────────────────────────────────────────────────
slide1(prs)                                   # 1  Title
section_break(prs, "SYSTEM OVERVIEW", 2)      # 2
slide3(prs)                                   # 3  What it does
slide4(prs)                                   # 4  Repo structure
slide5(prs)                                   # 5  Frontend stack
slide6(prs)                                   # 6  Backend stack
slide7(prs)                                   # 7  AI & External
section_break(prs, "BACKEND ARCHITECTURE", 8) # 8
slide9(prs)                                   # 9  NestJS modules
slide10(prs)                                  # 10 BullMQ
slide11(prs)                                  # 11 DB schema
slide12(prs)                                  # 12 Auth & security
section_break(prs, "FRONTEND ARCHITECTURE", 13) # 13
slide14(prs)                                  # 14 App Router
slide15(prs)                                  # 15 SSR vs client
slide16(prs)                                  # 16 Plan gating
section_break(prs, "THE AI AGENT SYSTEM", 17)   # 17
slide18(prs)                                  # 18 Agent principles
slide19(prs)                                  # 19 15 agents
slide20(prs)                                  # 20 New agent checklist
section_break(prs, "DEPLOYMENT & CI/CD", 21)    # 21
slide22(prs)                                  # 22 Deploy arch
slide23(prs)                                  # 23 Env vars
slide24(prs)                                  # 24 CI/CD pipeline
slide25(prs)                                  # 25 Testing
slide26(prs)                                  # 26 Golden rules
slide27(prs)                                  # 27 Common pitfalls
slide28(prs)                                  # 28 Closing

# ── Save ──────────────────────────────────────────────────────────────────────
output_path = r"D:\project\creatorforce-ai\CreatorForce-Developer-Guide.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slide count: {len(prs.slides)}")
