"""Rename all CreatorForce references to Sozialzync in .pptx files, then rename the files."""
import os
import re
from pptx import Presentation
from pptx.util import Pt
from copy import deepcopy

ROOT = r"D:\project\creatorforce-ai"

REPLACEMENTS = [
    ("AI CreatorForce", "Sozialzync"),
    ("AI-CreatorForce", "Sozialzync"),
    ("CreatorForce AI", "Sozialzync"),
    ("CreatorForce", "Sozialzync"),
    ("creatorforce", "sozialzync"),
    ("CREATORFORCE", "SOZIALZYNC"),
    ("AI Creatorforce", "Sozialzync"),
    ("Creatorforce", "Sozialzync"),
]

PPTX_RENAMES = {
    "AI-CreatorForce-Presentation.pptx": "Sozialzync-Complete-Presentation.pptx",
    "CreatorForce-Creator-Guide.pptx": "Sozialzync-Creator-Guide.pptx",
    "CreatorForce-Developer-Guide.pptx": "Sozialzync-Developer-Guide.pptx",
    "CreatorForce-Enterprise-Guide.pptx": "Sozialzync-Enterprise-Guide.pptx",
    "CreatorForce-SuperAdmin-Guide.pptx": "Sozialzync-SuperAdmin-Guide.pptx",
}

def replace_text(text):
    for old, new in REPLACEMENTS:
        text = text.replace(old, new)
    return text

def fix_run(run):
    if run.text:
        run.text = replace_text(run.text)

def fix_shape(shape):
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                fix_run(run)
    # Tables
    if shape.shape_type == 19:  # TABLE
        table = shape.table
        for row in table.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        fix_run(run)

def process_pptx(src_path, dst_path):
    prs = Presentation(src_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            fix_shape(shape)
            # Handle grouped shapes
            if shape.shape_type == 6:  # GROUP
                for s in shape.shapes:
                    fix_shape(s)
    prs.save(dst_path)
    print(f"  Saved: {os.path.basename(dst_path)}")

print("=" * 60)
print("  Renaming CreatorForce -> Sozialzync in all .pptx files")
print("=" * 60)

for old_name, new_name in PPTX_RENAMES.items():
    src = os.path.join(ROOT, old_name)
    dst = os.path.join(ROOT, new_name)
    if os.path.exists(src):
        print(f"\nProcessing: {old_name}")
        process_pptx(src, dst)
        os.remove(src)
        print(f"  Deleted old: {old_name}")
    else:
        print(f"  SKIP (not found): {old_name}")

# Also rename the python helper scripts
print("\nRenaming make_*.py scripts...")
for fname in os.listdir(ROOT):
    if fname.startswith("make_") and fname.endswith(".pptx.py") or \
       fname in ["make_pptx.py", "make_creator_pptx.py", "make_dev_pptx.py",
                 "make_superadmin_pptx.py", "make_enterprise_pptx.py",
                 "rename_to_sozialzync.py"]:
        new_fname = fname.replace("creatorforce", "sozialzync").replace("CreatorForce", "Sozialzync")
        if new_fname != fname:
            os.rename(os.path.join(ROOT, fname), os.path.join(ROOT, new_fname))
            print(f"  {fname} -> {new_fname}")

print("\nDone! Listing new .pptx files:")
for f in sorted(os.listdir(ROOT)):
    if f.endswith(".pptx"):
        size = os.path.getsize(os.path.join(ROOT, f))
        print(f"  {f}  ({size//1024} KB)")
