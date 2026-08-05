"""
AI CreatorForce — Creator Guide PowerPoint Generator
Generates a 22-slide presentation for YouTube content creators.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── Colour palette ────────────────────────────────────────────────────────────
PURPLE       = RGBColor(0x7C, 0x3A, 0xED)
ACCENT       = RGBColor(0xA7, 0x8B, 0xFA)
LIGHT_PURPLE = RGBColor(0xC4, 0xB5, 0xFD)
DARK_TEXT    = RGBColor(0x1E, 0x1B, 0x4B)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG     = RGBColor(0xF5, 0xF3, 0xFF)
GRAY         = RGBColor(0x9C, 0xA3, 0xAF)
ALT_ROW      = RGBColor(0xF5, 0xF3, 0xFF)

# ── Slide dimensions ──────────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

# ── Layout constants ──────────────────────────────────────────────────────────
HEADER_H     = Inches(1.2)
BODY_Y       = Inches(1.4)
BODY_H       = Inches(5.8)
BODY_X       = Inches(0.4)
BODY_W       = Inches(12.4)


# ─────────────────────────────────────────────────────────────────────────────
# Helper utilities
# ─────────────────────────────────────────────────────────────────────────────

def blank_slide(prs):
    """Return a new slide using the blank layout."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_rect(slide, x, y, w, h, fill_rgb, line=False):
    """Add a filled rectangle; optionally remove its border."""
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if not line:
        shape.line.fill.background()
    return shape


def add_textbox(slide, x, y, w, h, text, font_size, color, bold=False,
                align=PP_ALIGN.LEFT, italic=False, word_wrap=True):
    """Add a text box with a single run."""
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = word_wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold  = bold
    run.font.italic = italic
    return txb


def add_slide_number(slide, num):
    """Add slide number in bottom-right corner."""
    add_textbox(slide,
                Inches(12.0), Inches(7.1), Inches(1.1), Inches(0.3),
                str(num), 10, GRAY, align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────────────────────────────────────
# Slide builders
# ─────────────────────────────────────────────────────────────────────────────

def add_header_bar(slide, title_text):
    """Add purple header bar + white title text inside it."""
    add_rect(slide, Inches(0), Inches(0), W, HEADER_H, PURPLE)
    add_textbox(slide,
                Inches(0.2), Inches(0.25),
                Inches(12.8), Inches(0.85),
                title_text, 28, WHITE, bold=True)


def make_full_purple_slide(prs, title, subtitle=None, slide_num=None,
                            title_size=44, sub_size=18):
    """Full-purple background section-break slide."""
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, W, H, PURPLE)

    title_y = Inches(2.8) if subtitle else Inches(3.1)
    add_textbox(slide,
                Inches(0.5), title_y,
                Inches(12.3), Inches(1.2),
                title, title_size, WHITE, bold=True, align=PP_ALIGN.CENTER)

    if subtitle:
        add_textbox(slide,
                    Inches(0.5), Inches(4.2),
                    Inches(12.3), Inches(0.7),
                    subtitle, sub_size, LIGHT_PURPLE, align=PP_ALIGN.CENTER)

    if slide_num:
        add_textbox(slide,
                    Inches(12.0), Inches(7.1), Inches(1.1), Inches(0.3),
                    str(slide_num), 10, LIGHT_PURPLE, align=PP_ALIGN.RIGHT)
    return slide


def make_bullet_slide(prs, title, bullets, slide_num):
    """Header-bar slide with a bullet list."""
    slide = blank_slide(prs)
    add_rect(slide, Inches(0), Inches(0), W, H, LIGHT_BG)   # subtle bg
    add_header_bar(slide, title)

    # Build multi-paragraph text box for bullets
    txb = slide.shapes.add_textbox(BODY_X, BODY_Y, BODY_W, BODY_H)
    tf  = txb.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # line spacing via XML
        pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', '130000')   # 130 %

        run = p.add_run()
        run.text = bullet
        run.font.size = Pt(15)
        run.font.color.rgb = DARK_TEXT

    add_slide_number(slide, slide_num)
    return slide


# ─────────────────────────────────────────────────────────────────────────────
# Table helper (slide 20)
# ─────────────────────────────────────────────────────────────────────────────

def set_cell_fill(cell, rgb):
    """Set table cell background fill via lxml."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    # Remove existing solidFill if present
    for old in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(old)
    solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', f'{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}')


def make_table_slide(prs, slide_num):
    """Plan comparison table slide."""
    slide = blank_slide(prs)
    add_rect(slide, Inches(0), Inches(0), W, H, LIGHT_BG)
    add_header_bar(slide, "Plan Comparison")

    rows = 5   # header + 4 data rows
    cols = 4
    left = Inches(0.5)
    top  = Inches(1.55)
    tw   = Inches(12.3)
    th   = Inches(4.8)

    tbl_shape = slide.shapes.add_table(rows, cols, left, top, tw, th)
    tbl = tbl_shape.table

    # Column widths
    tbl.columns[0].width = Inches(4.2)
    tbl.columns[1].width = Inches(2.5)
    tbl.columns[2].width = Inches(2.8)
    tbl.columns[3].width = Inches(2.8)

    headers = ["Feature", "Free", "Pro (Credits)", "Enterprise"]
    data = [
        ["Projects",                "3",        "Unlimited",         "Unlimited"],
        ["AI Copilot queries/day",  "10",       "Unlimited",         "Unlimited"],
        ["Shorts edits/month",      "10",       "Unlimited",         "Unlimited"],
        ["AI Agents",               "Basic",    "All 15",            "All 15 + Priority"],
    ]

    # Header row
    for c, hdr in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = hdr
        set_cell_fill(cell, (0x7C, 0x3A, 0xED))
        for para in cell.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                run.font.size  = Pt(14)
                run.font.bold  = True
                run.font.color.rgb = WHITE

    # Data rows
    for r, row_data in enumerate(data):
        bg = (0xF5, 0xF3, 0xFF) if r % 2 == 0 else (0xFF, 0xFF, 0xFF)
        for c, val in enumerate(row_data):
            cell = tbl.cell(r + 1, c)
            cell.text = val
            set_cell_fill(cell, bg)
            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in para.runs:
                    run.font.size = Pt(13)
                    run.font.color.rgb = DARK_TEXT

    add_slide_number(slide, slide_num)
    return slide


# ─────────────────────────────────────────────────────────────────────────────
# Slide 1 — Title
# ─────────────────────────────────────────────────────────────────────────────

def make_title_slide(prs):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, W, H, PURPLE)

    add_textbox(slide,
                Inches(0.5), Inches(2.5), Inches(12.3), Inches(0.9),
                "AI CreatorForce", 52, WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_textbox(slide,
                Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.6),
                "Your AI-Powered YouTube Content Team",
                24, WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide,
                Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.5),
                "Research. Script. Comply. Publish. Grow.",
                16, LIGHT_PURPLE, align=PP_ALIGN.CENTER)

    add_textbox(slide,
                Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.4),
                "sozialzync.vercel.app",
                14, WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide,
                Inches(12.0), Inches(7.1), Inches(1.1), Inches(0.3),
                "1", 10, LIGHT_PURPLE, align=PP_ALIGN.RIGHT)
    return slide


# ─────────────────────────────────────────────────────────────────────────────
# Slide 22 — CTA
# ─────────────────────────────────────────────────────────────────────────────

def make_cta_slide(prs):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, W, H, PURPLE)

    add_textbox(slide,
                Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.8),
                "Ready to grow your channel?",
                36, WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_textbox(slide,
                Inches(0.5), Inches(3.0), Inches(12.3), Inches(0.7),
                "sozialzync.vercel.app",
                32, WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_textbox(slide,
                Inches(0.5), Inches(3.9), Inches(12.3), Inches(0.55),
                "Sign up free — no credit card required",
                20, WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide,
                Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.4),
                "AI CreatorForce — Your YouTube Content Operating System",
                14, LIGHT_PURPLE, align=PP_ALIGN.CENTER)

    add_textbox(slide,
                Inches(12.0), Inches(7.1), Inches(1.1), Inches(0.3),
                "22", 10, LIGHT_PURPLE, align=PP_ALIGN.RIGHT)
    return slide


# ─────────────────────────────────────────────────────────────────────────────
# Main build
# ─────────────────────────────────────────────────────────────────────────────

def build_presentation():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    make_title_slide(prs)

    # ── Slide 2: Section break ────────────────────────────────────────────────
    make_full_purple_slide(prs, "FOR THE CONTENT CREATOR",
                           subtitle=None, slide_num=2, title_size=44)

    # ── Slide 3: What is AI CreatorForce? ────────────────────────────────────
    make_bullet_slide(prs, "What is AI CreatorForce?", [
        "● AI-powered YouTube Content Operating System",
        "● Replaces an entire content team with 15 specialized AI agents",
        "● From idea to published video — fully guided workflow",
        "● Designed for monetizable, original, copyright-safe content",
        "● Works alongside you — you always approve before anything publishes",
        "● Supports Video, YouTube Shorts, and Music content types",
    ], 3)

    # ── Slide 4: Getting Started ──────────────────────────────────────────────
    make_bullet_slide(prs, "Getting Started — Your First 5 Minutes", [
        "● Sign up at sozialzync.vercel.app (email OTP or Google login)",
        "● Create your first Project — name it, pick a content type",
        "● Connect your YouTube channel (YouTube Data API OAuth)",
        "● Set your niche and brand voice in Channel Profile",
        "● Let the AI suggest your first video topic",
    ], 4)

    # ── Slide 5: Free Plan ────────────────────────────────────────────────────
    make_bullet_slide(prs, "Your Free Plan — What You Get", [
        "● 3 Projects (channels)",
        "● Up to 5 AI-generated outputs per project",
        "● 10 AI Copilot queries per day",
        "● 10 Shorts edits per month",
        "● Basic Research + Script generation",
        "● Human approval always required before publishing",
    ], 5)

    # ── Slide 6: Credits & Pro ────────────────────────────────────────────────
    make_bullet_slide(prs, "Upgrading: Credits & Pro Access", [
        "● Pro access: buy AI credits (pay-as-you-go, no subscription needed)",
        "● Credits unlock: all 15 AI agents, unlimited projects, full SEO suite",
        "● Credit types: Trial (welcome bonus) → Purchased → Bonus → Referral",
        "● Credits never expire during active use (oldest lot consumed first)",
        "● Enterprise Plan: team seats, shared org wallet, white-label features",
        "● Manage credits at: App → Wallet → Buy Credits",
    ], 6)

    # ── Slide 7: Section break — Content Pipeline ─────────────────────────────
    make_full_purple_slide(prs, "THE CONTENT PIPELINE",
                           subtitle="From Idea to Published Video",
                           slide_num=7, title_size=40, sub_size=18)

    # ── Slide 8: Step 1 — Research ────────────────────────────────────────────
    make_bullet_slide(prs, "Step 1 — Research & Topic Discovery", [
        "● ResearchAgent scans YouTube trends, search volume, and competition",
        "● Finds low-competition, high-opportunity video topics in your niche",
        "● Analyzes top-performing competitor videos for gap opportunities",
        "● Outputs: suggested titles, estimated reach, difficulty score",
        "● AudienceAgent profiles your target viewers and their pain points",
        "● You review and pick the topic — AI doesn't decide for you",
    ], 8)

    # ── Slide 9: Step 2 — Script ──────────────────────────────────────────────
    make_bullet_slide(prs, "Step 2 — Script Writing", [
        "● ScriptAgent generates a full script: hook → body → CTA",
        "● Hooks are engineered for audience retention (first 30 seconds critical)",
        "● Every factual claim flagged for verification before inclusion",
        "● Brand voice settings applied automatically (tone, style, vocabulary)",
        "● Script broken into scenes: talking head, B-roll cues, graphics notes",
        "● You edit, approve, or regenerate any section",
    ], 9)

    # ── Slide 10: Step 3 — Compliance ────────────────────────────────────────
    make_bullet_slide(prs, "Step 3 — Fact Checking & Compliance", [
        "● FactCheckAgent verifies every factual claim in the script",
        "● Each verified fact gets a source citation stored in the project",
        "● ComplianceAgent runs monetization safety analysis",
        "● Checks: advertiser-friendliness, copyright risk, community guidelines",
        "● Assigns a Monetization Safety Score (0–100)",
        "● Nothing moves to publishing until compliance passes — this is a hard gate",
    ], 10)

    # ── Slide 11: Step 4 — SEO ───────────────────────────────────────────────
    make_bullet_slide(prs, "Step 4 — SEO & Metadata Optimization", [
        "● SEOAgent generates: title variants, description, 30+ tags, chapters",
        "● Titles optimized for YouTube search ranking + click-through rate",
        "● Description includes natural keyword placement + chapter timestamps",
        "● MetadataAgent creates: cards, end-screen plan, pinned comment draft",
        "● ThumbnailAgent generates: detailed thumbnail creative brief",
        "● You choose from multiple title/thumbnail options",
    ], 11)

    # ── Slide 12: Step 5 — Review & Approve ──────────────────────────────────
    make_bullet_slide(prs, "Step 5 — Review & Approve", [
        "● All AI outputs land in your Project dashboard for review",
        "● Side-by-side comparison: AI draft vs. your notes",
        "● One-click regenerate for any section you dislike",
        "● Compliance score displayed prominently before you can publish",
        "● You must click \"Approve\" — publishing never happens automatically",
        "● Scheduled publish: you set the time, AI doesn't override it",
    ], 12)

    # ── Slide 13: Step 6 — Publish ───────────────────────────────────────────
    make_bullet_slide(prs, "Step 6 — Publish to YouTube", [
        "● Click Publish → system uploads video via YouTube Data API",
        "● Title, description, tags, chapters, cards all applied automatically",
        "● Thumbnail uploaded if you approved the AI-generated one",
        "● Video visibility: Public / Unlisted / Scheduled — you decide",
        "● Publishing log saved for compliance audit trail",
        "● Post-publish: AudienceAgent monitors early performance signals",
    ], 13)

    # ── Slide 14: Section break — Shorts Studio ───────────────────────────────
    make_full_purple_slide(prs, "SHORTS STUDIO",
                           subtitle="Turn Long Videos into Viral Shorts",
                           slide_num=14, title_size=40, sub_size=18)

    # ── Slide 15: Shorts Studio — Overview ───────────────────────────────────
    make_bullet_slide(prs, "Shorts Studio — Overview", [
        "● Import any of your existing YouTube videos into Shorts Studio",
        "● AI analyzes transcript + audio to find high-retention clip moments",
        "● Scene detection identifies natural cut points automatically",
        "● Generates up to 10 short clip candidates per video",
        "● Each clip: trimmed, reframed for vertical (9:16), captioned",
        "● Compliance check applied to every clip before export",
    ], 15)

    # ── Slide 16: Shorts Studio — Workflow ───────────────────────────────────
    make_bullet_slide(prs, "Shorts Studio — Workflow", [
        "● Import → Analyze → Clip Selection → Edit → Export → Publish",
        "● You select which clips to keep from AI recommendations",
        "● In-browser editor: trim, adjust captions, add music cue",
        "● Export as MP4 or publish directly to YouTube as a Short",
        "● Separate SEO metadata generated for each Short",
        "● Free plan: 10 Shorts edits/month; Pro: unlimited",
    ], 16)

    # ── Slide 17: AI Copilot ──────────────────────────────────────────────────
    make_bullet_slide(prs, "AI Copilot — Your Always-On Assistant", [
        "● Built-in chat assistant on every screen — ask anything",
        "● \"Write me a hook for a video about morning routines\"",
        "● \"What tags should I use for my cooking channel?\"",
        "● \"Is this title optimized for search?\"",
        "● Understands your channel profile and past content",
        "● Free: 10 queries/day | Pro: unlimited",
    ], 17)

    # ── Slide 18: Analytics ───────────────────────────────────────────────────
    make_bullet_slide(prs, "Analytics & Performance Tracking", [
        "● Dashboard shows: views, watch time, CTR, subscriber growth",
        "● Tracks which AI-generated videos perform best",
        "● Identifies content patterns that correlate with high retention",
        "● Weekly performance summary email (Pro feature)",
        "● AudienceAgent updates viewer profile based on engagement data",
        "● Recommendations: what to make next, best posting times",
    ], 18)

    # ── Slide 19: Section break — Plan Options ────────────────────────────────
    make_full_purple_slide(prs, "YOUR PLAN OPTIONS",
                           subtitle=None, slide_num=19, title_size=44)

    # ── Slide 20: Plan Comparison table ──────────────────────────────────────
    make_table_slide(prs, 20)

    # ── Slide 21: Tips for Success ────────────────────────────────────────────
    make_bullet_slide(prs, "Tips for Success on AI CreatorForce", [
        "● Set up your Channel Profile first — AI uses your brand voice everywhere",
        "● Always review the Compliance Score before publishing",
        "● Use Copilot to iterate on hooks — retention starts in second 1",
        "● Import your best existing video into Shorts Studio immediately",
        "● Buy a small credit pack to unlock all 15 agents for your first real video",
        "● Check analytics weekly — let data guide your next topic pick",
    ], 21)

    # ── Slide 22: CTA ─────────────────────────────────────────────────────────
    make_cta_slide(prs)

    out_path = r"D:\project\creatorforce-ai\CreatorForce-Creator-Guide.pptx"
    prs.save(out_path)
    print(f"Saved: {out_path}")
    print(f"Slides: {len(prs.slides)}")
    return out_path, len(prs.slides)


if __name__ == "__main__":
    path, count = build_presentation()
    print(f"\nDone. {count} slides written to:\n  {path}")
